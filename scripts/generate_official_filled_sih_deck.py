import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_official_filled_sih_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette strictly matching SIH Official Template
    BG_WHITE = RGBColor(255, 255, 255)
    NAVY_PRIMARY = RGBColor(29, 78, 137)     # #1D4E89 Official SIH Header Blue
    FOOTER_BLUE = RGBColor(0, 115, 198)      # #0073C6 Official SIH Footer Bar Blue
    TEXT_BLACK = RGBColor(0, 0, 0)           # Pure Black for titles and bullets
    TEXT_DARK = RGBColor(30, 41, 59)         # #1E293B Deep Slate for body
    TEXT_MUTED = RGBColor(71, 85, 105)       # #475569 Secondary
    
    PURPLE_OUTLINE = RGBColor(147, 51, 234)  # #9333EA Thin purple oval outline
    
    BORDER_BLUE = RGBColor(2, 132, 199)      # #0284C7 Accent card border
    BG_LIGHT_BLUE = RGBColor(240, 249, 255)  # #F0F9FF Light blue card
    
    BORDER_GREEN = RGBColor(22, 163, 74)     # #16A34A Green accent border
    BG_LIGHT_GREEN = RGBColor(240, 253, 244) # #F0FDF4 Light green card
    
    BORDER_RED = RGBColor(220, 38, 38)       # #DC2626 Red accent border
    BG_LIGHT_RED = RGBColor(254, 242, 242)   # #FEF2F2 Light red card
    
    BORDER_ORANGE = RGBColor(234, 88, 12)    # #EA580C Orange accent border
    BG_LIGHT_AMBER = RGBColor(254, 243, 199) # #FEF3C7 Light amber card
    
    BORDER_GREY = RGBColor(203, 213, 225)    # #CBD5E1 Neutral border
    BG_LIGHT_GREY = RGBColor(248, 250, 252)  # #F8FAFC Light grey card

    SIH_LOGO_PATH = "/tmp/sih_logo_extracted.png"
    HEX_GRAPHIC_PATH = "/tmp/sih_hex_graphic.png"
    DIAGRAM_S2_PATH = "/tmp/sih_assets/slide2_flowchart.png"
    DIAGRAM_S3_PATH = "/tmp/sih_assets/slide3_gis_curves.png"
    DIAGRAM_S5_PATH = "/tmp/sih_assets/slide5_timeline.png"

    def set_white_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_WHITE

    def add_official_header(slide, slide_title):
        # Top-Left Team Oval Badge (Thin purple oval with stacked lines)
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(0.18), Inches(1.5), Inches(0.75))
        oval.fill.solid()
        oval.fill.fore_color.rgb = BG_WHITE
        oval.line.color.rgb = PURPLE_OUTLINE
        oval.line.width = Pt(1.5)
        tf_ov = oval.text_frame
        tf_ov.word_wrap = False
        tf_ov.margin_left = tf_ov.margin_top = tf_ov.margin_right = tf_ov.margin_bottom = Inches(0)
        
        p1 = tf_ov.paragraphs[0]
        p1.text = "Team"
        p1.font.size = Pt(9)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_BLACK
        p1.font.name = "Arial"
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf_ov.add_paragraph()
        p2.text = "Thermo-"
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_BLACK
        p2.font.name = "Arial"
        p2.alignment = PP_ALIGN.CENTER
        
        p3 = tf_ov.add_paragraph()
        p3.text = "Shield"
        p3.font.size = Pt(9)
        p3.font.bold = True
        p3.font.color.rgb = TEXT_BLACK
        p3.font.name = "Arial"
        p3.alignment = PP_ALIGN.CENTER

        # Top-Center Slide Title (Serif Uppercase Bold)
        t_box = slide.shapes.add_textbox(Inches(2.4), Inches(0.2), Inches(8.3), Inches(0.65))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = Inches(0)
        p_t = tf_t.paragraphs[0]
        p_t.text = slide_title.upper()
        p_t.font.size = Pt(17)
        p_t.font.bold = True
        p_t.font.name = "Times New Roman"
        p_t.font.color.rgb = TEXT_BLACK
        p_t.alignment = PP_ALIGN.CENTER

        # Top-Right Official SIH Logo Image
        if os.path.exists(SIH_LOGO_PATH):
            slide.shapes.add_picture(SIH_LOGO_PATH, Inches(10.8), Inches(0.15), width=Inches(2.1))

    def add_official_footer(slide, slide_num):
        # Bottom Blue Footer Bar
        footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
        footer_bar.fill.solid()
        footer_bar.fill.fore_color.rgb = FOOTER_BLUE
        footer_bar.line.fill.background()

        # Footer Text
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(7.12), Inches(11.733), Inches(0.35))
        tf = tb.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = "@SIH Idea submission- Template"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

        # Slide Number on Right
        num_box = slide.shapes.add_textbox(Inches(12.2), Inches(7.12), Inches(0.8), Inches(0.35))
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.font.size = Pt(11)
        p_num.font.bold = True
        p_num.font.color.rgb = RGBColor(255, 255, 255)
        p_num.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 1: TITLE PAGE
    # =========================================================================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s1)

    # Top Header Text: SMART INDIA HACKATHON 2026
    h1_box = s1.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(9.8), Inches(0.6))
    tf_h1 = h1_box.text_frame
    p_h1 = tf_h1.paragraphs[0]
    p_h1.text = "SMART INDIA HACKATHON 2026"
    p_h1.font.size = Pt(26)
    p_h1.font.bold = True
    p_h1.font.name = "Times New Roman"
    p_h1.font.color.rgb = NAVY_PRIMARY

    # Top Right SIH Logo
    if os.path.exists(SIH_LOGO_PATH):
        s1.shapes.add_picture(SIH_LOGO_PATH, Inches(10.8), Inches(0.22), width=Inches(2.1))

    # Centered TITLE PAGE Label
    tp_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.733), Inches(0.5))
    tf_tp = tp_box.text_frame
    p_tp = tf_tp.paragraphs[0]
    p_tp.text = "TITLE PAGE"
    p_tp.font.size = Pt(16)
    p_tp.font.bold = True
    p_tp.font.name = "Times New Roman"
    p_tp.font.color.rgb = TEXT_BLACK
    p_tp.alignment = PP_ALIGN.CENTER

    # Right Background Hexagon Graphic
    if os.path.exists(HEX_GRAPHIC_PATH):
        s1.shapes.add_picture(HEX_GRAPHIC_PATH, Inches(7.8), Inches(1.7), width=Inches(4.9))

    # Left Column: Exact Template Bulleted Fields
    fields_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.8), Inches(4.5))
    tf_f = fields_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = Inches(0)

    title_fields = [
        ("Problem Statement ID –", "SIH26083"),
        ("Problem Statement Title-", "Extreme Heatwave Early Warning and Human Thermal Stress Index"),
        ("Theme-", "Disaster Management"),
        ("PS Category-", "Software"),
        ("Organization-", "Ministry of Earth Sciences (MoES) / NCMRWF"),
        ("Team ID-", "[YOUR REGISTERED TEAM ID]"),
        ("Team Name (Registered on portal)-", "ThermoShield")
    ]

    for idx, (lbl, val) in enumerate(title_fields):
        p = tf_f.paragraphs[0] if idx == 0 else tf_f.add_paragraph()
        p.text = f"• {lbl} "
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = TEXT_BLACK
        if idx > 0:
            p.space_before = Pt(7)
        
        run_v = p.add_run()
        run_v.text = val
        run_v.font.bold = False
        run_v.font.color.rgb = NAVY_PRIMARY if "ID" in lbl or "Name" in lbl else TEXT_BLACK

    # Left Bottom Solution Hero Callout
    sol_hero = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.65), Inches(6.6), Inches(1.3))
    sol_hero.fill.solid()
    sol_hero.fill.fore_color.rgb = BG_LIGHT_BLUE
    sol_hero.line.color.rgb = BORDER_BLUE
    sol_hero.line.width = Pt(1.5)
    tf_sh = sol_hero.text_frame
    tf_sh.word_wrap = True
    
    p_sh1 = tf_sh.paragraphs[0]
    p_sh1.text = "PROPOSED SYSTEM: ThermoShield"
    p_sh1.font.size = Pt(11.5)
    p_sh1.font.bold = True
    p_sh1.font.color.rgb = NAVY_PRIMARY
    
    p_sh2 = tf_sh.add_paragraph()
    p_sh2.text = "AI-Driven Human Heat Risk & Biometeorological Early Warning OS\nFrom Weather Forecasts → Physiological Thermal Stress → Hyper-Local Ward Action"
    p_sh2.font.size = Pt(9.5)
    p_sh2.font.color.rgb = TEXT_DARK
    p_sh2.space_before = Pt(2)

    # =========================================================================
    # SLIDE 2: IDEA TITLE (WITH EMBEDDED FLOWCHART DIAGRAM)
    # =========================================================================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s2)
    add_official_header(s2, "IDEA TITLE: ThermoShield — AI-Driven Human Heat Risk OS")

    # Section Header
    sec_tb = s2.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(11.733), Inches(0.35))
    tf_sec = sec_tb.text_frame
    tf_sec.margin_left = tf_sec.margin_top = tf_sec.margin_right = tf_sec.margin_bottom = Inches(0)
    p_sec = tf_sec.paragraphs[0]
    p_sec.text = "❖  Proposed Solution (Describe your Idea/Solution/Prototype)"
    p_sec.font.size = Pt(11)
    p_sec.font.bold = True
    p_sec.font.underline = True
    p_sec.font.color.rgb = NAVY_PRIMARY
    p_sec.font.name = "Arial"

    # Embedded Flowchart Architecture Image
    if os.path.exists(DIAGRAM_S2_PATH):
        s2.shapes.add_picture(DIAGRAM_S2_PATH, Inches(0.8), Inches(1.42), width=Inches(11.733), height=Inches(1.85))

    # 3 Structured Bottom Cards mapped to the 3 official template sub-bullets
    col_w2 = Inches(3.75)
    col_gap2 = Inches(0.24)
    left_start2 = Inches(0.8)

    cards_s2 = [
        ("• Detailed Explanation of the Solution", [
            ("Dual-Branch Ingestion:", "Fuses 5-Day hourly NWP forecasts (T, RH, Wind, Rad) with Census 2011 PCA demographics."),
            ("UTCI Physiological Engine:", "Universal Thermal Climate Index (6th-order polynomial based on Fiala 187-node model)."),
            ("ISO 7243 WBGT:", "Natural wet bulb (Stull) and globe temp (Liljegren) for occupational and outdoor labor stress."),
            ("Multi-Day Persistence:", "Exponential duration multiplier (Dmult) and uncooled night penalty (Tmin > 28°C).")
        ], BORDER_BLUE, BG_LIGHT_BLUE),
        ("• How It Addresses the Problem", [
            ("Human Impact vs Weather:", "Overcomes dry-bulb limitations where 40°C with high humidity breaks thermoregulation."),
            ("Hyper-Local Ward Granularity:", "Differentiates shaded affluent zones from informal settlements with tin roofs."),
            ("Epidemiological Grounding:", "Calibrated against Indian empirical studies (Azhar 2014 & Mazdiyasni 2017)."),
            ("Relative Risk Indexing:", "Bounded 0–100 Human Heat-Health Risk Score replacing uncalibrated death claims.")
        ], BORDER_GREEN, BG_LIGHT_GREEN),
        ("• Innovation and Uniqueness", [
            ("72h–120h Pre-Emptive Window:", "3–5 day anticipation enabling municipal water tankers and cooling shelters to be staged."),
            ("Automated Multi-Sector Triggers:", "Converts risk into hospital surge beds, DISCOM load warnings, and NIOSH labor halts."),
            ("Bilingual Public Advisories:", "Plain-language vernacular WhatsApp/Telegram warnings in Hindi and English."),
            ("Production Open Stack:", "25 REST APIs with sub-second latency and zero proprietary software licensing.")
        ], BORDER_ORANGE, BG_LIGHT_AMBER)
    ]

    for c_idx, (c_title, c_items, c_border, c_bg) in enumerate(cards_s2):
        c_left = left_start2 + c_idx * (col_w2 + col_gap2)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, Inches(3.35), col_w2, Inches(3.55))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = c_border
        card.line.width = Pt(1.5)

        # Header ribbon
        rib = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left + Inches(0.06), Inches(3.42), col_w2 - Inches(0.12), Inches(0.35))
        rib.fill.solid()
        rib.fill.fore_color.rgb = c_bg
        rib.line.color.rgb = c_border
        rib.line.width = Pt(1)
        tf_rib = rib.text_frame
        p_rib = tf_rib.paragraphs[0]
        p_rib.text = c_title
        p_rib.font.size = Pt(9.5)
        p_rib.font.bold = True
        p_rib.font.color.rgb = NAVY_PRIMARY
        p_rib.alignment = PP_ALIGN.CENTER

        # Body text
        tb = s2.shapes.add_textbox(c_left + Inches(0.12), Inches(3.85), col_w2 - Inches(0.24), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)

        for i_idx, (head, body) in enumerate(c_items):
            p = tf.paragraphs[0] if i_idx == 0 else tf.add_paragraph()
            p.text = f"❖  {head} "
            p.font.size = Pt(8.5)
            p.font.bold = True
            p.font.color.rgb = NAVY_PRIMARY
            if i_idx > 0:
                p.space_before = Pt(3)
            
            run_b = p.add_run()
            run_b.text = body
            run_b.font.bold = False
            run_b.font.color.rgb = TEXT_DARK

    add_official_footer(s2, 2)

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH (WITH GIS MOCKUP & FORECAST CURVES)
    # =========================================================================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s3)
    add_official_header(s3, "TECHNICAL APPROACH")

    # Section 1 Header
    m_tb = s3.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(11.733), Inches(0.3))
    tf_m = m_tb.text_frame
    tf_m.margin_left = tf_m.margin_top = tf_m.margin_right = tf_m.margin_bottom = Inches(0)
    p_m = tf_m.paragraphs[0]
    p_m.text = "• Methodology & Implementation Process (Working Prototype Architecture & Live GIS Output)"
    p_m.font.size = Pt(11)
    p_m.font.bold = True
    p_m.font.color.rgb = NAVY_PRIMARY
    p_m.font.name = "Arial"

    # Left: Process Pipeline Steps (Card)
    card_steps = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(3.5))
    card_steps.fill.solid()
    card_steps.fill.fore_color.rgb = BG_WHITE
    card_steps.line.color.rgb = BORDER_BLUE
    card_steps.line.width = Pt(1.5)

    tf_cs = card_steps.text_frame
    tf_cs.word_wrap = True
    tf_cs.margin_left = tf_cs.margin_top = tf_cs.margin_right = tf_cs.margin_bottom = Inches(0)

    p = tf_cs.paragraphs[0]
    p.text = "4-STAGE COMPUTATIONAL PIPELINE"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    steps_data = [
        ("1. Ingestion:", "Open-Meteo 0.1° NWP 5-day hourly weather + NASA POWER 40-yr baseline + Census 2011 PCA demographics."),
        ("2. Thermal Engine:", "UTCI 6th-order polynomial (Fiala) + ISO 7243 WBGT psychrometrics + Multi-day persistence penalty (Dmult)."),
        ("3. AI & HVI Synthesis:", "PCA Demographic Vulnerability Index (slums, elderly, gig labor, canopy deficit) → 0–100 Heat-Health Score."),
        ("4. Automated Action:", "Interactive Leaflet GIS map, 25 REST APIs, NDMA playbooks, hospital cooling beds, NIOSH labor halts.")
    ]
    for s_head, s_body in steps_data:
        p_s = tf_cs.add_paragraph()
        p_s.text = f"• {s_head} "
        p_s.font.size = Pt(8.5)
        p_s.font.bold = True
        p_s.font.color.rgb = NAVY_PRIMARY
        p_s.space_before = Pt(3)

        run_sb = p_s.add_run()
        run_sb.text = s_body
        run_sb.font.bold = False
        run_sb.font.color.rgb = TEXT_DARK

    # Right: Working Prototype Diagram (Choropleth + Curves)
    if os.path.exists(DIAGRAM_S3_PATH):
        s3.shapes.add_picture(DIAGRAM_S3_PATH, Inches(6.6), Inches(1.4), width=Inches(5.933), height=Inches(3.5))

    # Section 2: Technologies to be used (Bottom Categorized Container)
    tech_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.05), Inches(11.733), Inches(1.85))
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = BG_LIGHT_GREY
    tech_box.line.color.rgb = BORDER_GREY
    tech_box.line.width = Pt(1.2)
    tf_tb = tech_box.text_frame
    tf_tb.word_wrap = True

    p = tf_tb.paragraphs[0]
    p.text = "• Technologies to be used (Programming Languages, Frameworks, GIS & Deployment):"
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    tech_categories = [
        ("Backend & REST APIs:", "Python 3.11, FastAPI, Pydantic v2, Uvicorn, AsyncIO (25 REST production endpoints)"),
        ("Biometeorology & GIS Engine:", "NumPy, Pandas, GeoPandas, Shapely, Scipy (UTCI 6th-order polynomial & ISO 7243 psychrometrics)"),
        ("Web GIS & Visualization:", "Leaflet.js (Choropleth GIS), Chart.js (5-Day Trend Curves), HTML5/CSS3 Responsive UI"),
        ("Data Ingestion & Deployment:", "Open-Meteo NWP, NASA POWER MERRA-2, Census India PCA, Docker, Linux VPS (Oracle Linux)")
    ]

    for cat, desc in tech_categories:
        p_cat = tf_tb.add_paragraph()
        p_cat.text = f"❖  {cat} "
        p_cat.font.size = Pt(9)
        p_cat.font.bold = True
        p_cat.font.color.rgb = NAVY_PRIMARY
        p_cat.space_before = Pt(2)

        run_desc = p_cat.add_run()
        run_desc.text = desc
        run_desc.font.bold = False
        run_desc.font.color.rgb = TEXT_DARK

    add_official_footer(s3, 3)

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s4)
    add_official_header(s4, "FEASIBILITY AND VIABILITY")

    # Section 1: Feasibility Analysis Header & Narrative
    top_tb4 = s4.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(11.733), Inches(1.5))
    tf_t4 = top_tb4.text_frame
    tf_t4.word_wrap = True
    tf_t4.margin_left = tf_t4.margin_top = tf_t4.margin_right = tf_t4.margin_bottom = Inches(0)

    p1 = tf_t4.paragraphs[0]
    p1.text = "• Analysis of the Feasibility of the Idea"
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = NAVY_PRIMARY

    p1_sub = tf_t4.add_paragraph()
    p1_sub.text = "The ThermoShield platform is highly feasible due to open, keyless Tier-1 atmospheric data pipelines (Open-Meteo 0.1° NWP hourly forecasts and NASA POWER MERRA-2 40-year climatological normals) and validated biometeorological formulations (WMO UTCI polynomial & ISO 7243 WBGT). Integrating Census 2011 PCA demographics enables hyper-local ward-level vulnerability mapping without waiting for restricted government permissions."
    p1_sub.font.size = Pt(9)
    p1_sub.font.color.rgb = TEXT_DARK
    p1_sub.space_before = Pt(2)

    # Lower Left: Potential Challenges and Risks & Strategies for Overcoming
    c_left4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.65), Inches(7.5), Inches(4.3))
    c_left4.fill.solid()
    c_left4.fill.fore_color.rgb = BG_WHITE
    c_left4.line.color.rgb = BORDER_BLUE
    c_left4.line.width = Pt(1.5)

    tb_h4 = s4.shapes.add_textbox(Inches(1.0), Inches(2.75), Inches(7.1), Inches(0.35))
    tf_h4 = tb_h4.text_frame
    p_h4 = tf_h4.paragraphs[0]
    p_h4.text = "• Potential Challenges & Risks                     • Strategies for Overcoming"
    p_h4.font.size = Pt(10.5)
    p_h4.font.bold = True
    p_h4.font.color.rgb = BORDER_RED

    challenge_pairs = [
        ("Data Heterogeneity (0.1° Grid vs Wards)", "Spatial Harmonization & Polygon Attribution"),
        ("Confidential Hospital Mortality Records", "Epidemiological Calibration (Azhar 2014 & Mazdiyasni 2017)"),
        ("Model Overfitting / False Death Claims", "Bounded 0–100 Human Heat-Health Risk Score"),
        ("False Alarm Fatigue (Single Spikes)", "Compound Multi-Day Persistence Multiplier (Dmult)"),
        ("Uncooled Nighttime Strain (Tmin > 28°C)", "Explicit Nocturnal Recovery Penalty Algorithm")
    ]

    for idx, (ch, st) in enumerate(challenge_pairs):
        tb_pair = s4.shapes.add_textbox(Inches(1.0), Inches(3.18 + idx * 0.72), Inches(7.1), Inches(0.65))
        tf_p = tb_pair.text_frame
        tf_p.word_wrap = True
        tf_p.margin_left = tf_p.margin_top = tf_p.margin_right = tf_p.margin_bottom = Inches(0)
        p = tf_p.paragraphs[0]
        
        run_ch = p.add_run()
        run_ch.text = f"• {ch}\n"
        run_ch.font.size = Pt(9)
        run_ch.font.bold = True
        run_ch.font.color.rgb = BORDER_RED

        run_arr = p.add_run()
        run_arr.text = f"   ➔  {st}"
        run_arr.font.size = Pt(9)
        run_arr.font.bold = True
        run_arr.font.color.rgb = BORDER_GREEN

    # Lower Right: Expected Customers & Commercial Viability
    c_right4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(2.65), Inches(4.033), Inches(4.3))
    c_right4.fill.solid()
    c_right4.fill.fore_color.rgb = BG_LIGHT_GREY
    c_right4.line.color.rgb = NAVY_PRIMARY
    c_right4.line.width = Pt(1.5)

    tf_cr4 = c_right4.text_frame
    tf_cr4.word_wrap = True
    p_cr_head = tf_cr4.paragraphs[0]
    p_cr_head.text = "• Target Users & Viability:"
    p_cr_head.font.size = Pt(11)
    p_cr_head.font.bold = True
    p_cr_head.font.color.rgb = NAVY_PRIMARY

    target_users = [
        ("Municipal Corporations & Smart Cities", "For ward-level water tankers, cool roofs & cooling centers."),
        ("Health Departments & Hospitals", "For heatstroke emergency ward & IV fluid prepositioning."),
        ("State Disaster Management (SDMA / NDMA)", "For automated Heat Action Plan (HAP) trigger activation."),
        ("Gig Platforms & Construction Sector", "For mandatory NIOSH labor halt (11 AM–4 PM) compliance."),
        ("DISCOMs & Power Utilities", "For predictive grid load planning during peak cooling demand.")
    ]
    for c_title, c_desc in target_users:
        p_c = tf_cr4.add_paragraph()
        p_c.text = f"➢  {c_title}"
        p_c.font.size = Pt(9)
        p_c.font.bold = True
        p_c.font.color.rgb = NAVY_PRIMARY
        p_c.space_before = Pt(5)

        p_cd = tf_cr4.add_paragraph()
        p_cd.text = f"    {c_desc}"
        p_cd.font.size = Pt(8)
        p_cd.font.color.rgb = TEXT_MUTED

    add_official_footer(s4, 4)

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS (WITH EMBEDDED TIMELINE INFOGRAPHIC)
    # =========================================================================
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s5)
    add_official_header(s5, "IMPACT AND BENEFITS")

    # Section 1 Header
    aud_box = s5.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(11.733), Inches(0.3))
    tf_aud = aud_box.text_frame
    tf_aud.word_wrap = True
    tf_aud.margin_left = tf_aud.margin_top = tf_aud.margin_right = tf_aud.margin_bottom = Inches(0)
    
    p_aud = tf_aud.paragraphs[0]
    p_aud.text = "• Potential Impact on the Target Audience (72h–120h Pre-Emptive Intervention Timeline)"
    p_aud.font.size = Pt(11)
    p_aud.font.bold = True
    p_aud.font.color.rgb = NAVY_PRIMARY

    # Embedded Timeline Diagram
    if os.path.exists(DIAGRAM_S5_PATH):
        s5.shapes.add_picture(DIAGRAM_S5_PATH, Inches(0.8), Inches(1.4), width=Inches(11.733), height=Inches(1.8))

    # Section 2: Benefits of the solution (3 Bordered Columns)
    col_w5 = Inches(3.75)
    col_gap5 = Inches(0.24)
    left_start5 = Inches(0.8)

    benefits_data = [
        ("• Social Benefits", [
            ("Protects Vulnerable Citizens:", "Dramatically reduces heatstroke morbidity among elderly, infants, and chronic disease patients."),
            ("Protects Informal Laborers:", "Enforces humane work-rest cycles for delivery and construction workers."),
            ("Bilingual Public Alerts:", "Delivers plain-language vernacular guidance via SMS/WhatsApp to eliminate technical barriers.")
        ], BORDER_BLUE),
        ("• Economic Benefits", [
            ("Prevents Productivity Loss:", "Mitigates lost labor hours through scheduled morning/evening shifts ($3.8B/yr India heat loss)."),
            ("Optimized Relief Routing:", "Avoids wasteful blanket relief by targeting water tankers to high-vulnerability wards."),
            ("Reduces Emergency Care Costs:", "Pre-empts severe multi-organ failure and ICU hospitalization surges.")
        ], BORDER_GREEN),
        ("• Environmental & Urban Benefits", [
            ("Urban Heat Island (UHI) Mapping:", "Identifies micro-urban hot spots and green canopy deficits."),
            ("Sustainable Urban Planning:", "Guides long-term municipal cool roof installations and tree planting programs."),
            ("Effective Climate Resilience:", "Operationalizes national disaster mitigation at municipal ward scale.")
        ], BORDER_ORANGE)
    ]

    for c_idx, (b_title, b_items, b_color) in enumerate(benefits_data):
        c_left_pos = left_start5 + c_idx * (col_w5 + col_gap5)
        card_b = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left_pos, Inches(3.3), col_w5, Inches(3.65))
        card_b.fill.solid()
        card_b.fill.fore_color.rgb = BG_WHITE
        card_b.line.color.rgb = b_color
        card_b.line.width = Pt(1.5)

        tf_cb = card_b.text_frame
        tf_cb.word_wrap = True
        tf_cb.margin_left = tf_cb.margin_top = tf_cb.margin_right = tf_cb.margin_bottom = Inches(0)

        p = tf_cb.paragraphs[0]
        p.text = b_title
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = b_color

        for b_h, b_d in b_items:
            p_itm = tf_cb.add_paragraph()
            p_itm.text = f"▪  {b_h} "
            p_itm.font.size = Pt(8.5)
            p_itm.font.bold = True
            p_itm.font.color.rgb = NAVY_PRIMARY
            p_itm.space_before = Pt(3)

            run_bd = p_itm.add_run()
            run_bd.text = b_d
            run_bd.font.bold = False
            run_bd.font.color.rgb = TEXT_DARK

    add_official_footer(s5, 5)

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # =========================================================================
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s6)
    add_official_header(s6, "RESEARCH AND REFERENCES")

    # Section Header
    ref_head = s6.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.733), Inches(0.35))
    tf_rh = ref_head.text_frame
    tf_rh.margin_left = tf_rh.margin_top = tf_rh.margin_right = tf_rh.margin_bottom = Inches(0)
    p_rh = tf_rh.paragraphs[0]
    p_rh.text = "• Details / Links of the Reference and Research Work"
    p_rh.font.size = Pt(11.5)
    p_rh.font.bold = True
    p_rh.font.color.rgb = NAVY_PRIMARY

    # References Body Box
    ref_box = s6.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.733), Inches(4.7))
    tf_r6 = ref_box.text_frame
    tf_r6.word_wrap = True
    tf_r6.margin_left = tf_r6.margin_top = tf_r6.margin_right = tf_r6.margin_bottom = Inches(0)

    references_list = [
        ("Universal Thermal Climate Index (UTCI) Consortium (WMO / COST Action 730):", "\"Development of the Universal Thermal Climate Index (UTCI)\", Bröde, Fiala, Błażejczyk et al., Int J Biometeorol (2012). Multi-node thermophysiological polynomial."),
        ("Occupational Heat Stress Standards (ISO 7243:2017 & NIOSH):", "\"Criteria for a Recommended Standard: Occupational Exposure to Heat and Hot Environments\", CDC / NIOSH Publication No. 2016-106. Standard for WBGT thresholds & work-rest cycles."),
        ("Epidemiological Heat-Health Grounding in Indian Cities (Ahmedabad 2010):", "\"The 2010 Ahmedabad Heat Wave: Impact on Mortality and Strategy for Resilient Cities\", Azhar et al., PLoS ONE (2014). Documented 1,344 excess deaths (43.1% surge) validating humidity-heat risk."),
        ("50-Year Multi-Decadal Heat Impact in India (PNAS Reference):", "\"Increasing probability of mass-fatality heatwaves in India\", Mazdiyasni et al., Proceedings of the National Academy of Sciences (PNAS, 2017). Proving +0.5°C threshold triggers 146% rise in mass mortality events."),
        ("National Disaster & Health Policy Guidelines (NDMA & MoHFW NCDC):", "\"National Action Plan for Heat-Related Illnesses (NAP-HRI)\", National Programme on Climate Change and Human Health (NPCCHH, 2024) & NDMA National Heatwave Guidelines."),
        ("NWP & Satellite Climatology Portals:", "National Centre for Medium Range Weather Forecasting (NCMRWF) NWP Portal, IMD Geospatial Bulletins & NASA POWER MERRA-2 (LaRC) 40-Year Climatological Baseline."),
        ("Demographic Spatial Vulnerability:", "Census of India 2011 Primary Census Abstract (PCA) Ward Demographics & Open-Meteo High-Resolution Numerical Weather Prediction API.")
    ]

    for idx, (r_title, r_desc) in enumerate(references_list):
        p = tf_r6.paragraphs[0] if idx == 0 else tf_r6.add_paragraph()
        p.text = f"• {r_title} "
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = NAVY_PRIMARY
        if idx > 0:
            p.space_before = Pt(3.5)

        run_rd = p.add_run()
        run_rd.text = r_desc
        run_rd.font.bold = False
        run_rd.font.color.rgb = TEXT_DARK

    # Bottom Gap Addressed Box
    gap_box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.15), Inches(11.733), Inches(0.8))
    gap_box.fill.solid()
    gap_box.fill.fore_color.rgb = BG_LIGHT_BLUE
    gap_box.line.color.rgb = BORDER_BLUE
    gap_box.line.width = Pt(1.2)
    tf_gb = gap_box.text_frame
    tf_gb.word_wrap = True
    
    p_gb = tf_gb.paragraphs[0]
    p_gb.text = "RESEARCH GAP SOLVED: Existing systems provide either macro-weather forecasts OR static vulnerability maps. ThermoShield operationalizes the complete unified pipeline: Forecast → Physiological Thermal Stress → Demographic Vulnerability → Ward Risk → Automated Action."
    p_gb.font.size = Pt(9.5)
    p_gb.font.bold = True
    p_gb.font.color.rgb = NAVY_PRIMARY
    p_gb.alignment = PP_ALIGN.CENTER

    add_official_footer(s6, 6)

    # Save to canonical target locations
    out_path_1 = "/home/ubuntu/sih26083-heat-risk/docs/sih/SIH2026_IDEA_Presentation_SIH26083_Official.pptx"
    out_path_2 = "/home/ubuntu/sih26083_research/SIH26083_SIH_Official_Submission.pptx"

    os.makedirs(os.path.dirname(out_path_1), exist_ok=True)
    os.makedirs(os.path.dirname(out_path_2), exist_ok=True)

    prs.save(out_path_1)
    prs.save(out_path_2)
    print(f"Master SIH 2026 Presentation successfully generated and saved to:\n1. {out_path_1}\n2. {out_path_2}")

if __name__ == "__main__":
    build_official_filled_sih_deck()
