import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_sih_master_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette from Official Winning References (Ref 1-4)
    BG_WHITE = RGBColor(255, 255, 255)
    NAVY_PRIMARY = RGBColor(15, 44, 89)      # #0F2C59 Deep SIH Navy
    ORANGE_SIH = RGBColor(234, 88, 12)       # #EA580C SIH Saffron / Orange
    PURPLE_BADGE = RGBColor(109, 40, 217)    # #6D28D9 Team Oval Badge
    PURPLE_BG = RGBColor(243, 232, 255)      # #F3E8FF Team Badge Fill
    
    TEXT_DARK = RGBColor(15, 23, 42)         # #0F172A Primary Dark Text
    TEXT_MUTED = RGBColor(71, 85, 105)       # #475569 Subtitle / Secondary
    
    BORDER_BLUE = RGBColor(2, 132, 199)      # #0284C7 Blue Accent Border
    BG_LIGHT_BLUE = RGBColor(240, 249, 255)  # #F0F9FF Light Blue Card Fill
    
    BORDER_GREEN = RGBColor(22, 163, 74)     # #16A34A Green Accent Border
    BG_LIGHT_GREEN = RGBColor(240, 253, 244) # #F0FDF4 Light Green Card Fill
    
    BORDER_RED = RGBColor(220, 38, 38)       # #DC2626 Red Accent Border
    BG_LIGHT_RED = RGBColor(254, 242, 242)   # #FEF2F2 Light Red Card Fill
    
    BORDER_AMBER = RGBColor(217, 119, 6)     # #D97706 Amber Accent Border
    BG_LIGHT_AMBER = RGBColor(254, 243, 199) # #FEF3C7 Light Amber Card Fill
    
    BORDER_GREY = RGBColor(203, 213, 225)    # #CBD5E1 Neutral Border
    BG_LIGHT_GREY = RGBColor(248, 250, 252)  # #F8FAFC Light Grey Card Fill

    def set_white_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_WHITE

    def add_reference_header(slide, slide_title, team_name="ThermoShield"):
        # Top-Left Team Badge (Oval matching References 1, 2, 3, 4)
        team_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.2), Inches(2.2), Inches(0.55))
        team_shape.fill.solid()
        team_shape.fill.fore_color.rgb = PURPLE_BG
        team_shape.line.color.rgb = PURPLE_BADGE
        team_shape.line.width = Pt(1.5)
        tf_team = team_shape.text_frame
        tf_team.word_wrap = False
        p_tm = tf_team.paragraphs[0]
        p_tm.text = team_name
        p_tm.font.size = Pt(11)
        p_tm.font.bold = True
        p_tm.font.color.rgb = PURPLE_BADGE
        p_tm.font.name = "Arial"
        p_tm.alignment = PP_ALIGN.CENTER

        # Top-Center Slide Title (Serif Heading)
        t_box = slide.shapes.add_textbox(Inches(2.9), Inches(0.18), Inches(7.5), Inches(0.6))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = Inches(0)
        p_t = tf_t.paragraphs[0]
        p_t.text = slide_title
        p_t.font.size = Pt(17)
        p_t.font.bold = True
        p_t.font.name = "Georgia"
        p_t.font.color.rgb = TEXT_DARK
        p_t.alignment = PP_ALIGN.CENTER

        # Top-Right SIH 2026 Official Logo Emblem
        sih_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.15), Inches(2.2), Inches(0.65))
        sih_shape.fill.solid()
        sih_shape.fill.fore_color.rgb = BG_LIGHT_AMBER
        sih_shape.line.color.rgb = ORANGE_SIH
        sih_shape.line.width = Pt(1.5)
        tf_sih = sih_shape.text_frame
        p_s1 = tf_sih.paragraphs[0]
        p_s1.text = "SMART INDIA"
        p_s1.font.size = Pt(9.5)
        p_s1.font.bold = True
        p_s1.font.color.rgb = NAVY_PRIMARY
        p_s1.font.name = "Arial"
        p_s1.alignment = PP_ALIGN.CENTER
        p_s2 = tf_sih.add_paragraph()
        p_s2.text = "HACKATHON 2026"
        p_s2.font.size = Pt(9.5)
        p_s2.font.bold = True
        p_s2.font.color.rgb = ORANGE_SIH
        p_s2.font.name = "Arial"
        p_s2.alignment = PP_ALIGN.CENTER

    def add_reference_footer(slide, slide_num):
        # Bottom Blue Footer Bar
        footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
        footer_bar.fill.solid()
        footer_bar.fill.fore_color.rgb = NAVY_PRIMARY
        footer_bar.line.fill.background()

        # Footer Text
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(7.12), Inches(11.733), Inches(0.35))
        tf = tb.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = "@SIH Idea submission- Template"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

        # Slide Number
        num_box = slide.shapes.add_textbox(Inches(12.2), Inches(7.12), Inches(0.8), Inches(0.35))
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.font.size = Pt(11)
        p_num.font.bold = True
        p_num.font.color.rgb = RGBColor(255, 255, 255)
        p_num.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 1: TITLE PAGE (OFFICIAL FORMAT FROM REF 1, 2, 3, 4)
    # =========================================================================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s1)

    # Top Header
    h1_box = s1.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(9.5), Inches(0.6))
    tf_h1 = h1_box.text_frame
    p_h1 = tf_h1.paragraphs[0]
    p_h1.text = "SMART INDIA HACKATHON 2026"
    p_h1.font.size = Pt(26)
    p_h1.font.bold = True
    p_h1.font.name = "Georgia"
    p_h1.font.color.rgb = NAVY_PRIMARY

    # Top Right SIH Logo
    sih_logo_s1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.3), Inches(2.2), Inches(0.8))
    sih_logo_s1.fill.solid()
    sih_logo_s1.fill.fore_color.rgb = BG_LIGHT_AMBER
    sih_logo_s1.line.color.rgb = ORANGE_SIH
    sih_logo_s1.line.width = Pt(1.5)
    tf_ls1 = sih_logo_s1.text_frame
    p_ls1 = tf_ls1.paragraphs[0]
    p_ls1.text = "SMART INDIA"
    p_ls1.font.size = Pt(11)
    p_ls1.font.bold = True
    p_ls1.font.color.rgb = NAVY_PRIMARY
    p_ls1.alignment = PP_ALIGN.CENTER
    p_ls2 = tf_ls1.add_paragraph()
    p_ls2.text = "HACKATHON 2026"
    p_ls2.font.size = Pt(11)
    p_ls2.font.bold = True
    p_ls2.font.color.rgb = ORANGE_SIH
    p_ls2.alignment = PP_ALIGN.CENTER

    # Project Title (Large Serif Title Case)
    title_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.733), Inches(0.8))
    tf_title = title_box.text_frame
    p_t = tf_title.paragraphs[0]
    p_t.text = "ThermoShield — AI-Driven Human Heat Risk & Early Warning System"
    p_t.font.size = Pt(22)
    p_t.font.bold = True
    p_t.font.name = "Georgia"
    p_t.font.color.rgb = TEXT_DARK

    # Left Column: Structured Metadata Fields with Underlined Labels
    details_box = s1.shapes.add_textbox(Inches(0.8), Inches(2.05), Inches(6.8), Inches(4.8))
    tf_d = details_box.text_frame
    tf_d.word_wrap = True
    tf_d.margin_left = tf_d.margin_top = tf_d.margin_right = tf_d.margin_bottom = Inches(0)

    fields = [
        ("Problem Statement ID –", "SIH26083"),
        ("Problem Statement Title –", "Extreme Heatwave Early Warning and Human Thermal Stress Index"),
        ("Theme –", "Disaster Management"),
        ("PS Category –", "Software"),
        ("Organization –", "Ministry of Earth Sciences (MoES) / NCMRWF"),
        ("Team ID –", "[YOUR REGISTERED TEAM ID]"),
        ("Team Name –", "[YOUR REGISTERED TEAM NAME]")
    ]

    for idx, (label, val) in enumerate(fields):
        p = tf_d.paragraphs[0] if idx == 0 else tf_d.add_paragraph()
        p.text = f"• {label} "
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.name = "Arial"
        p.font.underline = True
        p.font.color.rgb = TEXT_DARK
        if idx > 0:
            p.space_before = Pt(7)
        
        run_v = p.add_run()
        run_v.text = val
        run_v.font.bold = False
        run_v.font.underline = False
        run_v.font.color.rgb = NAVY_PRIMARY if "ID" in label or "Name" in label else TEXT_DARK

    # Right Column: Visual Scope & Mission Infographic Card
    card_r1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(2.05), Inches(4.9), Inches(4.7))
    card_r1.fill.solid()
    card_r1.fill.fore_color.rgb = BG_LIGHT_BLUE
    card_r1.line.color.rgb = BORDER_BLUE
    card_r1.line.width = Pt(1.5)

    tf_r1 = card_r1.text_frame
    tf_r1.word_wrap = True
    p_r1_head = tf_r1.paragraphs[0]
    p_r1_head.text = "CORE SYSTEM CAPABILITIES"
    p_r1_head.font.size = Pt(12)
    p_r1_head.font.bold = True
    p_r1_head.font.color.rgb = NAVY_PRIMARY
    p_r1_head.alignment = PP_ALIGN.CENTER

    scope_points = [
        ("From Weather to Health:", "Translates raw atmospheric metrics into physiological human thermal stress."),
        ("Multi-Metric Engine:", "Universal Thermal Climate Index (UTCI) + ISO 7243 WBGT + NOAA Heat Index."),
        ("Hyper-Local Ward GIS:", "Census 2011 PCA demographic overlay (slums, elderly, gig workers)."),
        ("Automated Triggers:", "Bilingual NDMA playbooks, hospital cooling surge alerts & NIOSH labor halts."),
        ("3–5 Day Lead Time:", "Enables 72h–120h pre-emptive municipal water and emergency resource staging.")
    ]
    for sp_head, sp_body in scope_points:
        p_sp = tf_r1.add_paragraph()
        p_sp.text = f"✓  {sp_head} "
        p_sp.font.size = Pt(10)
        p_sp.font.bold = True
        p_sp.font.color.rgb = NAVY_PRIMARY
        p_sp.space_before = Pt(6)
        
        run_sp = p_sp.add_run()
        run_sp.text = sp_body
        run_sp.font.bold = False
        run_sp.font.color.rgb = TEXT_DARK

    add_reference_footer(s1, 1)

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION (REF 2/3 THREE-COLUMN & CARD ARCHITECTURE)
    # =========================================================================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s2)
    add_reference_header(s2, "ThermoShield : AI-Driven Human Heat Risk & Early Warning System", "ThermoShield")

    # 3 Vertical Section Columns (Matching Ref 3 style: Detailed Explanation | How It Addresses | Key Outcomes)
    col_w2 = Inches(3.75)
    col_gap2 = Inches(0.24)
    left_start2 = Inches(0.8)

    sections_s2 = [
        ("1. Proposed Solution Architecture", [
            ("Dual-Branch Ingestion:", "Simultaneously processes 5-Day NWP forecasts (temperature, humidity, wind, radiation) and Census demographic vulnerability layers."),
            ("UTCI Physiological Engine:", "6th-order polynomial based on Fiala 187-node human body model. Accurately calculates real sweat evaporation efficiency."),
            ("ISO 7243 WBGT:", "Natural wet bulb (Stull) and globe temperature (Liljegren) for occupational and outdoor worker heat stress."),
            ("Multi-Day Persistence Penalty:", "Applies exponential multiplier (Dmult) for prolonged heatwaves and uncooled night minimums (Tmin > 28°C).")
        ], BORDER_BLUE, BG_LIGHT_BLUE),
        ("2. How It Addresses The Problem?", [
            ("Weather vs Physiological Impact:", "Overcomes the critical limitation of dry-bulb warnings where 40°C with high humidity turns lethal."),
            ("Hyper-Local Ward Stratification:", "Differentiates shaded affluent residential zones from densely built informal settlements with tin roofs."),
            ("Epidemiological Grounding:", "Risk scoring calibrated against peer-reviewed Indian empirical studies (Azhar et al. 2014 & Mazdiyasni et al. 2017)."),
            ("Automated Impact-to-Action:", "Translates raw risk directly into automated NDMA advisories, hospital beds, and labor work-rest cycles.")
        ], BORDER_GREEN, BG_LIGHT_GREEN),
        ("3. Key Differentiators & Impact", [
            ("3–5 Day Pre-Emptive Window:", "Provides 72h–120h lead time for municipal water tanker routing and cooling shelter deployment."),
            ("Sector-Specific Automation:", "Delivers tailored feeds for gig delivery platforms (Zomato/Swiggy), construction, and DISCOMs."),
            ("Bilingual Vernacular Alerts:", "Automated Hindi & English public alerts sent via WhatsApp/Telegram to avoid technical jargon."),
            ("High-Performance Open Stack:", "25 production REST API endpoints with sub-second response times and zero licensing fees.")
        ], ORANGE_SIH, BG_LIGHT_AMBER)
    ]

    for c_idx, (sec_title, sec_items, sec_border, sec_bg) in enumerate(sections_s2):
        c_left = left_start2 + c_idx * (col_w2 + col_gap2)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, Inches(1.05), col_w2, Inches(5.35))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = sec_border
        card.line.width = Pt(1.5)

        # Header Ribbon inside card
        ribbon = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left + Inches(0.08), Inches(1.15), col_w2 - Inches(0.16), Inches(0.45))
        ribbon.fill.solid()
        ribbon.fill.fore_color.rgb = sec_bg
        ribbon.line.color.rgb = sec_border
        ribbon.line.width = Pt(1)
        tf_rib = ribbon.text_frame
        p_rib = tf_rib.paragraphs[0]
        p_rib.text = sec_title
        p_rib.font.size = Pt(10.5)
        p_rib.font.bold = True
        p_rib.font.color.rgb = NAVY_PRIMARY
        p_rib.alignment = PP_ALIGN.CENTER

        # Body text inside card
        tb = s2.shapes.add_textbox(c_left + Inches(0.15), Inches(1.7), col_w2 - Inches(0.3), Inches(4.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)

        for i_idx, (head, body) in enumerate(sec_items):
            p = tf.paragraphs[0] if i_idx == 0 else tf.add_paragraph()
            p.text = f"❖  {head} "
            p.font.size = Pt(9.5)
            p.font.bold = True
            p.font.color.rgb = NAVY_PRIMARY
            if i_idx > 0:
                p.space_before = Pt(4)
            
            run_b = p.add_run()
            run_b.text = body
            run_b.font.bold = False
            run_b.font.color.rgb = TEXT_DARK

    # Bottom USP Callout
    usp_card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.45), Inches(11.733), Inches(0.5))
    usp_card.fill.solid()
    usp_card.fill.fore_color.rgb = BG_LIGHT_BLUE
    usp_card.line.color.rgb = BORDER_BLUE
    usp_card.line.width = Pt(1.2)
    tf_u = usp_card.text_frame
    p_u = tf_u.paragraphs[0]
    p_u.text = "KEY USP: Not another static heat map — an automated, end-to-end impact-to-action intelligence OS for heatwave resilience."
    p_u.font.size = Pt(10)
    p_u.font.bold = True
    p_u.font.color.rgb = NAVY_PRIMARY
    p_u.alignment = PP_ALIGN.CENTER

    add_reference_footer(s2, 2)

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH (REF 1/2 PROCESS FLOW & CATEGORIZED TECH STACK)
    # =========================================================================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s3)
    add_reference_header(s3, "TECHNICAL APPROACH", "ThermoShield")

    # Sub-heading banner
    lbl_s3 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.8), Inches(0.95), Inches(7.733), Inches(0.35))
    lbl_s3.fill.solid()
    lbl_s3.fill.fore_color.rgb = BG_LIGHT_BLUE
    lbl_s3.line.color.rgb = BORDER_BLUE
    lbl_s3.line.width = Pt(1)
    tf_ls3 = lbl_s3.text_frame
    p_ls3 = tf_ls3.paragraphs[0]
    p_ls3.text = "End-to-End Biometeorological & Action Pipeline : Implementation Process Flow"
    p_ls3.font.size = Pt(10)
    p_ls3.font.bold = True
    p_ls3.font.color.rgb = NAVY_PRIMARY
    p_ls3.alignment = PP_ALIGN.CENTER

    # 4 Flowchart Pipeline Containers
    flow_steps = [
        ("1. INGESTION", "Open-Meteo 0.1° NWP 5-Day Forecasts\n• Hourly T, RH, Wind, Radiation\n• NASA POWER 40-Yr Normal\n• Census 2011 Ward PCA Data\n• Satellite LST / NDVI Surface Grid", ORANGE_SIH, BG_LIGHT_AMBER),
        ("2. THERMAL ENGINE", "UTCI 6th-Order Polynomial\n• Fiala 187-Node Human Model\n• ISO 7243 WBGT Psychrometrics\n• Stull Tw & Liljegren Tg Models\n• Cumulative Multi-Day Dmult", BORDER_GREEN, BG_LIGHT_GREEN),
        ("3. AI & HVI MODEL", "Census Demographic Vulnerability\n• Elderly Density & Slum Ratio\n• Outdoor Labor Concentration\n• Tree Canopy Deficit Index\n• Relative Heat-Health Score", BORDER_RED, BG_LIGHT_RED),
        ("4. ACTION & GIS", "Interactive Leaflet GIS Map\n• 25 REST Production APIs\n• Automated NDMA Playbooks\n• Hospital Surge Bed Triggers\n• NIOSH Labor Halt Alerts", NAVY_PRIMARY, BG_LIGHT_BLUE)
    ]

    c_w3 = Inches(2.78)
    c_gap3 = Inches(0.2)
    left_start3 = Inches(0.8)

    for idx, (f_title, f_body, f_border, f_bg) in enumerate(flow_steps):
        c_left = left_start3 + idx * (c_w3 + c_gap3)
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, Inches(1.4), c_w3, Inches(3.7))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_WHITE
        box.line.color.rgb = f_border
        box.line.width = Pt(1.5)

        # Header tag
        tag = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left + Inches(0.08), Inches(1.48), c_w3 - Inches(0.16), Inches(0.4))
        tag.fill.solid()
        tag.fill.fore_color.rgb = f_bg
        tag.line.color.rgb = f_border
        tag.line.width = Pt(1)
        tf_tag = tag.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = f_title
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = f_border
        p_tag.alignment = PP_ALIGN.CENTER

        # Body
        tb = s3.shapes.add_textbox(c_left + Inches(0.12), Inches(1.95), c_w3 - Inches(0.24), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        lines = f_body.split("\n")
        for l_idx, line in enumerate(lines):
            p = tf.paragraphs[0] if l_idx == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT_DARK
            if l_idx == 0:
                p.font.bold = True
                p.font.color.rgb = NAVY_PRIMARY
            else:
                p.space_before = Pt(3)

    # Bottom Container: Technologies to be used (Categorized Grid matching Ref 1 & 2)
    tech_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(11.733), Inches(1.75))
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = BG_LIGHT_GREY
    tech_box.line.color.rgb = BORDER_GREY
    tech_box.line.width = Pt(1.2)
    tf_tb = tech_box.text_frame
    tf_tb.word_wrap = True

    p = tf_tb.paragraphs[0]
    p.text = "Technologies to be used:"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    tech_categories = [
        ("• Backend & REST APIs:", "Python 3.11, FastAPI, Pydantic v2, Uvicorn, AsyncIO (25 REST endpoints)"),
        ("• Biometeorological & GIS Engine:", "NumPy, Pandas, GeoPandas, Shapely, Scipy (UTCI 6th-order polynomial & psychrometrics)"),
        ("• Web GIS & Visualization:", "Leaflet.js (Choropleth GIS), Chart.js (5-Day Trend Curves), HTML5/CSS3 Responsive UI"),
        ("• Data Ingestion & Cloud Infra:", "Open-Meteo NWP, NASA POWER (MERRA-2), Census India PCA, Docker, Linux VPS (Oracle Linux)")
    ]

    for cat, desc in tech_categories:
        p_cat = tf_tb.add_paragraph()
        p_cat.text = f"{cat} "
        p_cat.font.size = Pt(9)
        p_cat.font.bold = True
        p_cat.font.color.rgb = NAVY_PRIMARY
        p_cat.space_before = Pt(2)

        run_desc = p_cat.add_run()
        run_desc.text = desc
        run_desc.font.bold = False
        run_desc.font.color.rgb = TEXT_DARK

    add_reference_footer(s3, 3)

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY (REF 1/4 DUAL CONTAINER + RISK MATRIX)
    # =========================================================================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s4)
    add_reference_header(s4, "FEASIBILITY AND VIABILITY", "ThermoShield")

    # Top Section: Feasibility Narrative & Product Offering
    top_tb4 = s4.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(11.733), Inches(1.6))
    tf_t4 = top_tb4.text_frame
    tf_t4.word_wrap = True
    tf_t4.margin_left = tf_t4.margin_top = tf_t4.margin_right = tf_t4.margin_bottom = Inches(0)

    p1 = tf_t4.paragraphs[0]
    p1.text = "❖  Is this Idea feasible ?"
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = NAVY_PRIMARY

    p1_sub = tf_t4.add_paragraph()
    p1_sub.text = "The proposed ThermoShield system is fully feasible due to open, keyless Tier-1 atmospheric data pipelines (Open-Meteo 0.1° NWP hourly forecasts and NASA POWER MERRA-2 40-year climatological normals) and validated biometeorological equations (UTCI polynomial & ISO 7243 WBGT). The use of Census 2011 PCA demographics enables hyper-local ward-level vulnerability mapping without waiting for restricted government clearance."
    p1_sub.font.size = Pt(9)
    p1_sub.font.color.rgb = TEXT_DARK
    p1_sub.space_before = Pt(1)

    p2 = tf_t4.add_paragraph()
    p2.text = "❖  Product Offering: We offer this platform to Municipal Corporations, State Disaster Management Authorities (SDMAs), Health Departments, and Gig Platforms as an automated, impact-to-action early warning and resource coordination OS."
    p2.font.size = Pt(9)
    p2.font.bold = True
    p2.font.color.rgb = NAVY_PRIMARY
    p2.space_before = Pt(3)

    # Left Container: Challenges & Risks vs Strategies for Overcoming (Ref 1/4 Blue Border)
    c_left4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.7), Inches(7.5), Inches(4.25))
    c_left4.fill.solid()
    c_left4.fill.fore_color.rgb = BG_WHITE
    c_left4.line.color.rgb = BORDER_BLUE
    c_left4.line.width = Pt(1.5)

    # Container Header
    tb_h4 = s4.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(7.1), Inches(0.4))
    tf_h4 = tb_h4.text_frame
    p_h4 = tf_h4.paragraphs[0]
    p_h4.text = "Challenges & Risks:                                             Strategies for Overcoming:"
    p_h4.font.size = Pt(10.5)
    p_h4.font.bold = True
    p_h4.font.color.rgb = BORDER_RED

    challenge_pairs = [
        ("Data Heterogeneity (0.1° Grid vs Wards)", "Spatial Harmonization & Polygon Attribution"),
        ("Confidential Hospital Mortality Records", "Epidemiological Calibration (Azhar 2014 & Mazdiyasni 2017)"),
        ("Model Overfitting / False Death Claims", "Bounded 0–100 Human Heat-Health Risk Score"),
        ("False Alarm Fatigue (Single Spikes)", "Compound Multi-Day Persistence Multiplier (Dmult)"),
        ("Uncooled Nighttime Strain (Tmin > 28°C)", "Explicit Nocturnal Recovery Penalty Algorithm")
    ]

    for idx, (ch, st) in enumerate(challenge_pairs):
        tb_pair = s4.shapes.add_textbox(Inches(1.0), Inches(3.25 + idx * 0.72), Inches(7.1), Inches(0.65))
        tf_p = tb_pair.text_frame
        tf_p.word_wrap = True
        tf_p.margin_left = tf_p.margin_top = tf_p.margin_right = tf_p.margin_bottom = Inches(0)
        p = tf_p.paragraphs[0]
        
        run_ch = p.add_run()
        run_ch.text = f"• {ch}\n"
        run_ch.font.size = Pt(9)
        run_ch.font.bold = True
        run_ch.font.color.rgb = BORDER_RED

        run_arr = p.add_run()
        run_arr.text = f"   ➔  {st}"
        run_arr.font.size = Pt(9)
        run_arr.font.bold = True
        run_arr.font.color.rgb = BORDER_GREEN

    # Right Container: Expected Customers & Users (Ref 1/4 Navy Border)
    c_right4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(2.7), Inches(4.033), Inches(4.25))
    c_right4.fill.solid()
    c_right4.fill.fore_color.rgb = BG_LIGHT_GREY
    c_right4.line.color.rgb = NAVY_PRIMARY
    c_right4.line.width = Pt(1.5)

    tf_cr4 = c_right4.text_frame
    tf_cr4.word_wrap = True
    p_cr_head = tf_cr4.paragraphs[0]
    p_cr_head.text = "Expected Customers & Users:"
    p_cr_head.font.size = Pt(11)
    p_cr_head.font.bold = True
    p_cr_head.font.color.rgb = NAVY_PRIMARY

    target_users = [
        ("Municipal Corporations & Smart Cities", "For ward-level water tankers, cool roofs & cooling centers."),
        ("Health Departments & Hospitals", "For heatstroke emergency ward & IV fluid prepositioning."),
        ("State Disaster Management (SDMA / NDMA)", "For automated Heat Action Plan (HAP) trigger activation."),
        ("Gig Platforms & Construction Sector", "For mandatory NIOSH labor halt (11 AM–4 PM) compliance."),
        ("DISCOMs & Power Utilities", "For predictive grid load planning during peak cooling demand.")
    ]
    for c_title, c_desc in target_users:
        p_c = tf_cr4.add_paragraph()
        p_c.text = f"➢  {c_title}"
        p_c.font.size = Pt(9)
        p_c.font.bold = True
        p_c.font.color.rgb = NAVY_PRIMARY
        p_c.space_before = Pt(5)

        p_cd = tf_cr4.add_paragraph()
        p_cd.text = f"    {c_desc}"
        p_cd.font.size = Pt(8)
        p_cd.font.color.rgb = TEXT_MUTED

    add_reference_footer(s4, 4)

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS (REF 1/4 THREE-COLUMN BENEFIT MATRIX)
    # =========================================================================
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s5)
    add_reference_header(s5, "IMPACT AND BENEFITS", "ThermoShield")

    # Top Target Audience Line
    aud_box = s5.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(11.733), Inches(0.75))
    tf_aud = aud_box.text_frame
    tf_aud.word_wrap = True
    tf_aud.margin_left = tf_aud.margin_top = tf_aud.margin_right = tf_aud.margin_bottom = Inches(0)
    
    p_aud = tf_aud.paragraphs[0]
    p_aud.text = "Target Audience — Municipal Commissioners, Chief Medical Officers, Disaster Response Teams (NDMA/SDRF), Gig Workers, Construction Laborers, Urban Slum Residents, General Public."
    p_aud.font.size = Pt(10)
    p_aud.font.bold = True
    p_aud.font.color.rgb = NAVY_PRIMARY

    p_aud_sub = tf_aud.add_paragraph()
    p_aud_sub.text = "(All above stakeholders receive tailored biometeorological intelligence and automated sector-specific action protocols.)"
    p_aud_sub.font.size = Pt(8.5)
    p_aud_sub.font.color.rgb = TEXT_MUTED
    p_aud_sub.space_before = Pt(2)

    # Middle Banner: 72h–120h Lead-Time Timeline
    time_shape = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(11.733), Inches(1.4))
    time_shape.fill.solid()
    time_shape.fill.fore_color.rgb = BG_LIGHT_BLUE
    time_shape.line.color.rgb = BORDER_BLUE
    time_shape.line.width = Pt(1.2)
    tf_tm = time_shape.text_frame
    tf_tm.word_wrap = True
    
    p_tm_h = tf_tm.paragraphs[0]
    p_tm_h.text = "72h–120h (3–5 DAY) PRE-EMPTIVE HEAT ACTION TIMELINE:"
    p_tm_h.font.size = Pt(10)
    p_tm_h.font.bold = True
    p_tm_h.font.color.rgb = NAVY_PRIMARY

    p_tm_b = tf_tm.add_paragraph()
    p_tm_b.text = "• 120h (D-5): NWP model flags thermal moisture surge → Early municipal advisory & DISCOM grid warnings issued\n• 72h (D-3): Cumulative heat persistence detected → Water tankers and mobile misting units dispatched to high-risk slums\n• 48h (D-2): Hospital surge beds & ice-bath cooling centers staged; emergency ORS supplies prepositioned\n• 24h (D-1): Mandatory NIOSH labor halts mandated for construction & delivery workers during peak hours (11 AM–4 PM)"
    p_tm_b.font.size = Pt(8.5)
    p_tm_b.font.color.rgb = TEXT_DARK
    p_tm_b.space_before = Pt(2)

    # Three-Column Benefits Containers (Social, Economical, Environmental)
    col_w5 = Inches(3.75)
    col_gap5 = Inches(0.24)
    left_start5 = Inches(0.8)

    benefits_data = [
        ("● Social Benefits:", [
            ("Protects Vulnerable Citizens:", "Dramatically reduces heatstroke morbidity among elderly, children, and chronic disease patients."),
            ("Protects Informal Laborers:", "Enforces humane work-rest cycles for delivery and construction workers."),
            ("Bilingual Public Alerts:", "Delivers plain-language vernacular guidance via SMS/WhatsApp to eliminate technical barriers.")
        ], BORDER_BLUE),
        ("● Economical Benefits:", [
            ("Prevents Productivity Loss:", "Mitigates lost labor hours through scheduled morning/evening shifts ($3.8B/yr India heat loss)."),
            ("Optimized Relief Routing:", "Avoids wasteful blanket relief by targeting water tankers to high-vulnerability wards."),
            ("Reduces Emergency Care Costs:", "Pre-empts severe multi-organ failure and ICU hospitalization surges.")
        ], BORDER_GREEN),
        ("● Environmental & Urban Benefits:", [
            ("Urban Heat Island (UHI) Mapping:", "Identifies micro-urban hot spots and green canopy deficits."),
            ("Sustainable Urban Planning:", "Guides long-term municipal cool roof installations and tree planting programs."),
            ("Effective Climate Resilience:", "Operationalizes national disaster mitigation at municipal ward scale.")
        ], ORANGE_SIH)
    ]

    for c_idx, (b_title, b_items, b_color) in enumerate(benefits_data):
        c_left_pos = left_start5 + c_idx * (col_w5 + col_gap5)
        card_b = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left_pos, Inches(3.25), col_w5, Inches(3.7))
        card_b.fill.solid()
        card_b.fill.fore_color.rgb = BG_WHITE
        card_b.line.color.rgb = b_color
        card_b.line.width = Pt(1.5)

        tf_cb = card_b.text_frame
        tf_cb.word_wrap = True
        tf_cb.margin_left = tf_cb.margin_top = tf_cb.margin_right = tf_cb.margin_bottom = Inches(0)

        p = tf_cb.paragraphs[0]
        p.text = b_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = b_color

        for b_h, b_d in b_items:
            p_itm = tf_cb.add_paragraph()
            p_itm.text = f"▪  {b_h} "
            p_itm.font.size = Pt(8.5)
            p_itm.font.bold = True
            p_itm.font.color.rgb = NAVY_PRIMARY
            p_itm.space_before = Pt(4)

            run_bd = p_itm.add_run()
            run_bd.text = b_d
            run_bd.font.bold = False
            run_bd.font.color.rgb = TEXT_DARK

    add_reference_footer(s5, 5)

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES (EXACT WINNING TEMPLATE FORMAT)
    # =========================================================================
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s6)
    add_reference_header(s6, "RESEARCH AND REFERENCES", "ThermoShield")

    # Content Box with ● bullets matching reference style
    ref_box = s6.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.733), Inches(5.9))
    tf_r6 = ref_box.text_frame
    tf_r6.word_wrap = True
    tf_r6.margin_left = tf_r6.margin_top = tf_r6.margin_right = tf_r6.margin_bottom = Inches(0)

    references_list = [
        ("Universal Thermal Climate Index (UTCI) Consortium (WMO / COST Action 730):", "\"Development of the Universal Thermal Climate Index (UTCI)\", Bröde, Fiala, Błażejczyk et al., Int J Biometeorol (2012). Multi-node thermophysiological polynomial."),
        ("Occupational Heat Stress Standards (ISO 7243:2017 & NIOSH):", "\"Criteria for a Recommended Standard: Occupational Exposure to Heat and Hot Environments\", CDC / NIOSH Publication No. 2016-106. Standard for WBGT thresholds & work-rest cycles."),
        ("Epidemiological Heat-Health Grounding in Indian Cities (Ahmedabad 2010):", "\"The 2010 Ahmedabad Heat Wave: Impact on Mortality and Strategy for Resilient Cities\", Azhar et al., PLoS ONE (2014). Documented 1,344 excess deaths (43.1% surge) validating humidity-heat risk."),
        ("50-Year Multi-Decadal Heat Impact in India (PNAS Reference):", "\"Increasing probability of mass-fatality heatwaves in India\", Mazdiyasni et al., Proceedings of the National Academy of Sciences (PNAS, 2017). Proving +0.5°C threshold triggers 146% rise in mass mortality events."),
        ("National Disaster & Health Policy Guidelines (NDMA & MoHFW NCDC):", "\"National Action Plan for Heat-Related Illnesses (NAP-HRI)\", National Programme on Climate Change and Human Health (NPCCHH, 2024) & NDMA National Heatwave Guidelines."),
        ("NWP & Satellite Climatology Portals:", "National Centre for Medium Range Weather Forecasting (NCMRWF) NWP Portal, IMD Geospatial Bulletins & NASA POWER MERRA-2 (LaRC) 40-Year Climatological Baseline."),
        ("Demographic Spatial Vulnerability:", "Census of India 2011 Primary Census Abstract (PCA) Ward Demographics & Open-Meteo High-Resolution Numerical Weather Prediction API."),
        ("Research Gap Addressed:", "Bridges the operational gap between macro-scale meteorological forecasts and actionable, localized public health interventions by providing an automated impact-to-action intelligence OS.")
    ]

    for idx, (r_title, r_desc) in enumerate(references_list):
        p = tf_r6.paragraphs[0] if idx == 0 else tf_r6.add_paragraph()
        p.text = f"●  {r_title} "
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = NAVY_PRIMARY
        p.space_before = Pt(3.5)

        run_rd = p.add_run()
        run_rd.text = r_desc
        run_rd.font.bold = False
        run_rd.font.color.rgb = TEXT_DARK

    add_reference_footer(s6, 6)

    # Save to official repository and research directories
    out_path_1 = "/home/ubuntu/sih26083-heat-risk/docs/sih/SIH2026_IDEA_Presentation_SIH26083_Official.pptx"
    out_path_2 = "/home/ubuntu/sih26083_research/SIH26083_SIH_Official_Submission.pptx"

    os.makedirs(os.path.dirname(out_path_1), exist_ok=True)
    os.makedirs(os.path.dirname(out_path_2), exist_ok=True)

    prs.save(out_path_1)
    prs.save(out_path_2)
    print(f"Master SIH 2026 Presentation successfully generated and saved to:\n1. {out_path_1}\n2. {out_path_2}")

if __name__ == "__main__":
    build_sih_master_deck()
