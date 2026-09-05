import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_sih_exact_reference_template_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette matching the reference presentation
    BG_WHITE = RGBColor(255, 255, 255)
    NAVY_SIH = RGBColor(15, 44, 89)         # #0F2C59 Deep SIH Blue
    ORANGE_SIH = RGBColor(234, 88, 12)      # #EA580C SIH Saffron / Orange
    PURPLE_BADGE = RGBColor(109, 40, 217)   # #6D28D9 Team Name Badge
    TEXT_BLACK = RGBColor(15, 23, 42)       # #0F172A Body text
    TEXT_MUTED = RGBColor(71, 85, 105)      # #475569 Secondary text
    BORDER_BLUE = RGBColor(2, 132, 199)     # #0284C7 Accent container border
    BORDER_RED = RGBColor(220, 38, 38)       # #DC2626 Warning / risk border
    BORDER_GREEN = RGBColor(22, 163, 74)    # #16A34A Success border
    BORDER_GREY = RGBColor(203, 213, 225)   # #CBD5E1 Neutral border
    BG_LIGHT_BLUE = RGBColor(240, 249, 255) # Light blue box fill
    BG_LIGHT_GRAY = RGBColor(248, 250, 252) # Light gray box fill

    def set_white_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_WHITE

    def add_reference_header(slide, slide_title, team_name="ThermoShield"):
        # Top-Left Team Badge (Oval/Rounded)
        team_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.25), Inches(2.2), Inches(0.55))
        team_shape.fill.solid()
        team_shape.fill.fore_color.rgb = RGBColor(243, 232, 255)
        team_shape.line.color.rgb = PURPLE_BADGE
        team_shape.line.width = Pt(1.5)
        tf_team = team_shape.text_frame
        tf_team.word_wrap = False
        p_tm = tf_team.paragraphs[0]
        p_tm.text = team_name
        p_tm.font.size = Pt(11)
        p_tm.font.bold = True
        p_tm.font.color.rgb = PURPLE_BADGE
        p_tm.alignment = PP_ALIGN.CENTER

        # Top-Center Slide Title
        t_box = slide.shapes.add_textbox(Inches(3.0), Inches(0.2), Inches(7.333), Inches(0.65))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = Inches(0)
        p_t = tf_t.paragraphs[0]
        p_t.text = slide_title
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.name = "Georgia"
        p_t.font.color.rgb = TEXT_BLACK
        p_t.alignment = PP_ALIGN.CENTER

        # Top-Right SIH 2026 Logo Badge
        sih_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.2), Inches(2.2), Inches(0.65))
        sih_shape.fill.solid()
        sih_shape.fill.fore_color.rgb = RGBColor(254, 243, 199)
        sih_shape.line.color.rgb = ORANGE_SIH
        sih_shape.line.width = Pt(1.5)
        tf_sih = sih_shape.text_frame
        p_s1 = tf_sih.paragraphs[0]
        p_s1.text = "SMART INDIA"
        p_s1.font.size = Pt(9.5)
        p_s1.font.bold = True
        p_s1.font.color.rgb = NAVY_SIH
        p_s1.alignment = PP_ALIGN.CENTER
        p_s2 = tf_sih.add_paragraph()
        p_s2.text = "HACKATHON 2026"
        p_s2.font.size = Pt(9.5)
        p_s2.font.bold = True
        p_s2.font.color.rgb = ORANGE_SIH
        p_s2.alignment = PP_ALIGN.CENTER

    def add_reference_footer(slide, slide_num):
        # Bottom Blue Footer Bar
        footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
        footer_bar.fill.solid()
        footer_bar.fill.fore_color.rgb = NAVY_SIH
        footer_bar.line.fill.background()

        # Footer Text
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(7.12), Inches(11.733), Inches(0.35))
        tf = tb.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = "@SIH Idea submission- Template"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

        # Slide Number on right
        num_box = slide.shapes.add_textbox(Inches(12.2), Inches(7.12), Inches(0.8), Inches(0.35))
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.font.size = Pt(11)
        p_num.font.bold = True
        p_num.font.color.rgb = RGBColor(255, 255, 255)
        p_num.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 1: TITLE SLIDE (EXACT REFERENCE FORMAT)
    # =========================================================================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s1)

    # Top Header
    h1_box = s1.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(9.5), Inches(0.6))
    tf_h1 = h1_box.text_frame
    p_h1 = tf_h1.paragraphs[0]
    p_h1.text = "SMART INDIA HACKATHON 2026"
    p_h1.font.size = Pt(26)
    p_h1.font.bold = True
    p_h1.font.name = "Georgia"
    p_h1.font.color.rgb = NAVY_SIH

    # Top Right SIH Logo
    sih_logo_s1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.35), Inches(2.2), Inches(0.8))
    sih_logo_s1.fill.solid()
    sih_logo_s1.fill.fore_color.rgb = RGBColor(254, 243, 199)
    sih_logo_s1.line.color.rgb = ORANGE_SIH
    sih_logo_s1.line.width = Pt(1.5)
    tf_ls1 = sih_logo_s1.text_frame
    p_ls1 = tf_ls1.paragraphs[0]
    p_ls1.text = "SMART INDIA"
    p_ls1.font.size = Pt(11)
    p_ls1.font.bold = True
    p_ls1.font.color.rgb = NAVY_SIH
    p_ls1.alignment = PP_ALIGN.CENTER
    p_ls2 = tf_ls1.add_paragraph()
    p_ls2.text = "HACKATHON 2026"
    p_ls2.font.size = Pt(11)
    p_ls2.font.bold = True
    p_ls2.font.color.rgb = ORANGE_SIH
    p_ls2.alignment = PP_ALIGN.CENTER

    # Project Title (Large Serif)
    title_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.8))
    tf_title = title_box.text_frame
    p_t = tf_title.paragraphs[0]
    p_t.text = "ThermoShield — AI-Driven Human Heat Risk & Early Warning System"
    p_t.font.size = Pt(22)
    p_t.font.bold = True
    p_t.font.name = "Georgia"
    p_t.font.color.rgb = TEXT_BLACK

    # Left Column: Problem Statement & Team Details
    details_box = s1.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(6.8), Inches(4.7))
    tf_d = details_box.text_frame
    tf_d.word_wrap = True
    tf_d.margin_left = tf_d.margin_top = tf_d.margin_right = tf_d.margin_bottom = Inches(0)

    fields = [
        ("Problem Statement ID –", "SIH26083"),
        ("Problem Statement Title –", "Extreme Heatwave Early Warning and Human Thermal Stress Index"),
        ("Theme –", "Disaster Management"),
        ("PS Category –", "Software"),
        ("Organization –", "Ministry of Earth Sciences (MoES) / NCMRWF"),
        ("Team ID –", "[YOUR REGISTERED TEAM ID]"),
        ("Team Name –", "[YOUR REGISTERED TEAM NAME]")
    ]

    for idx, (label, val) in enumerate(fields):
        p = tf_d.paragraphs[0] if idx == 0 else tf_d.add_paragraph()
        p.text = f"• {label} "
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = "Arial"
        p.font.underline = True
        p.font.color.rgb = TEXT_BLACK
        if idx > 0:
            p.space_before = Pt(8)
        
        run_v = p.add_run()
        run_v.text = val
        run_v.font.bold = False
        run_v.font.underline = False
        run_v.font.color.rgb = NAVY_SIH if "ID" in label or "Name" in label else TEXT_BLACK

    # Right Column: Visual Infographic Card
    card_r1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(2.1), Inches(4.9), Inches(4.6))
    card_r1.fill.solid()
    card_r1.fill.fore_color.rgb = BG_LIGHT_BLUE
    card_r1.line.color.rgb = BORDER_BLUE
    card_r1.line.width = Pt(1.5)

    tf_r1 = card_r1.text_frame
    tf_r1.word_wrap = True
    p_r1_head = tf_r1.paragraphs[0]
    p_r1_head.text = "CORE MISSION & SCOPE"
    p_r1_head.font.size = Pt(13)
    p_r1_head.font.bold = True
    p_r1_head.font.color.rgb = NAVY_SIH
    p_r1_head.alignment = PP_ALIGN.CENTER

    scope_points = [
        "From Weather Forecasts → Human Thermal Stress → Localized Action",
        "Universal Thermal Climate Index (UTCI) + ISO 7243 WBGT Core",
        "Census 2011 PCA Demographic Vulnerability Mapping (HVI)",
        "Automated Bilingual NDMA Action Playbooks & Hospital Triggers",
        "3–5 Day (72h–120h) Pre-Emptive Forecast & Intervention Engine"
    ]
    for sp in scope_points:
        p_sp = tf_r1.add_paragraph()
        p_sp.text = f"✓  {sp}"
        p_sp.font.size = Pt(10.5)
        p_sp.font.bold = True
        p_sp.font.color.rgb = TEXT_BLACK
        p_sp.space_before = Pt(8)

    add_reference_footer(s1, 1)

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION & IDEA TITLE (EXACT REFERENCE FORMAT)
    # =========================================================================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s2)
    add_reference_header(s2, "ThermoShield : AI-Driven Human Heat Risk & Early Warning System", "ThermoShield")

    # Content Body Box with ❖ bullets matching the reference style
    body_box = s2.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.733), Inches(5.7))
    tf_b2 = body_box.text_frame
    tf_b2.word_wrap = True
    tf_b2.margin_left = tf_b2.margin_top = tf_b2.margin_right = tf_b2.margin_bottom = Inches(0)

    s2_points = [
        ("Dual-Branch Biometeorological Model:", "System utilizes two parallel analytical pipelines: one for Atmospheric Dynamics (NWP temperature, humidity, wind velocity, solar radiation) and one for Demographic Vulnerability (Census 2011 PCA elderly density, slum households, outdoor labor concentration, canopy deficit). Mathematical algorithms extract physiological heat stress indicators while GIS spatial models calculate ward-level vulnerability."),
        ("Universal Thermal Stress Engine:", "Thermal stress is computed using the Universal Thermal Climate Index (UTCI, 6th-order polynomial based on the Fiala 187-node human thermoregulation model) and ISO 7243 Outdoor WBGT (Stull wet-bulb + Liljegren radiation model). Accurately models human sweat evaporation breakdown where 40°C with high humidity turns lethal."),
        ("Multi-Day Heat Persistence & Nighttime Penalty:", "Calculates cumulative thermal accumulation via exponential persistence duration multiplier (Dmult). Accounts for uncooled nocturnal minimum temperatures (Tmin > 28°C) that prevent physiological recovery and double cardiovascular stress."),
        ("AI Health Risk Model & Multi-Factor Loss Optimization:", "The model synthesizes atmospheric hazard with local socio-economic vulnerability into a normalized 0–100 Human Heat-Health Risk Score:", [
            ("Relative Health Risk Loss:", "Calibrated against empirical epidemiological benchmarks (Azhar et al. 2014 & Mazdiyasni et al. 2017)."),
            ("Monotonicity Constraint:", "Guarantees mathematically that higher humidity/temperature strictly escalates risk scores."),
            ("Ward Spatial Allocation Loss:", "Ensures hyper-local resource allocation is prioritized toward high-vulnerability informal settlements.")
        ]),
        ("Automated Action Engine & Multi-Sector Dispatch:", "Converts risk calculations directly into sector-specific automated triggers: Municipal water tanker routing, emergency hospital cooling ward surge beds, gig-worker NIOSH work-rest cycles (labor halts 11 AM–4 PM), and bilingual NDMA SMS/WhatsApp broadcasts."),
        ("Final Output & Evaluation Metrics:", "Delivers real-time interactive Leaflet GIS choropleths, 25 high-speed REST APIs, and automated multi-channel alert dispatch with validated sub-second execution.")
    ]

    for idx, item in enumerate(s2_points):
        p = tf_b2.paragraphs[0] if idx == 0 else tf_b2.add_paragraph()
        if len(item) == 2:
            title, desc = item
            p.text = f"❖  {title} "
            p.font.size = Pt(9.5)
            p.font.bold = True
            p.font.color.rgb = NAVY_SIH
            p.space_before = Pt(4)
            
            run_d = p.add_run()
            run_d.text = desc
            run_d.font.bold = False
            run_d.font.color.rgb = TEXT_BLACK
        else:
            title, desc, sub_items = item
            p.text = f"❖  {title} "
            p.font.size = Pt(9.5)
            p.font.bold = True
            p.font.color.rgb = NAVY_SIH
            p.space_before = Pt(4)
            
            run_d = p.add_run()
            run_d.text = desc
            run_d.font.bold = False
            run_d.font.color.rgb = TEXT_BLACK

            for s_head, s_body in sub_items:
                p_sub = tf_b2.add_paragraph()
                p_sub.text = f"     ■  {s_head} "
                p_sub.font.size = Pt(9)
                p_sub.font.bold = True
                p_sub.font.color.rgb = ORANGE_SIH
                p_sub.space_before = Pt(2)
                
                run_sb = p_sub.add_run()
                run_sb.text = s_body
                run_sb.font.bold = False
                run_sb.font.color.rgb = TEXT_BLACK

    add_reference_footer(s2, 2)

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH (FLOWCHART & ARCHITECTURE)
    # =========================================================================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s3)
    add_reference_header(s3, "TECHNICAL APPROACH", "ThermoShield")

    # Lower Center Label
    lbl_shape = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(1.05), Inches(7.333), Inches(0.35))
    lbl_shape.fill.solid()
    lbl_shape.fill.fore_color.rgb = BG_LIGHT_BLUE
    lbl_shape.line.color.rgb = BORDER_BLUE
    lbl_shape.line.width = Pt(1)
    tf_lbl = lbl_shape.text_frame
    p_lbl = tf_lbl.paragraphs[0]
    p_lbl.text = "ThermoShield Biometeorological & Action Model: Implementation Architecture"
    p_lbl.font.size = Pt(10)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = NAVY_SIH
    p_lbl.alignment = PP_ALIGN.CENTER

    # Architecture Flowchart (Containers & Blocks)
    # 1. INPUT
    c_in = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.55), Inches(2.2), Inches(3.8))
    c_in.fill.solid()
    c_in.fill.fore_color.rgb = BG_LIGHT_BLUE
    c_in.line.color.rgb = ORANGE_SIH
    c_in.line.width = Pt(1.5)
    tf_in = c_in.text_frame
    p = tf_in.paragraphs[0]
    p.text = "1. MULTI-INPUT\nINGESTION"
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = ORANGE_SIH
    p.alignment = PP_ALIGN.CENTER
    
    in_items = [
        "NWP 5-Day Hourly Weather (0.1° Grid)",
        "NASA POWER 40-Yr Climatology Normal",
        "Census 2011 PCA Ward Demographics",
        "Satellite LST & NDVI Land Cover"
    ]
    for itm in in_items:
        p_itm = tf_in.add_paragraph()
        p_itm.text = f"• {itm}"
        p_itm.font.size = Pt(8.5)
        p_itm.font.color.rgb = TEXT_BLACK
        p_itm.space_before = Pt(4)

    # 2. PROCESSING & THERMAL STRESS
    c_proc = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(1.55), Inches(3.2), Inches(3.8))
    c_proc.fill.solid()
    c_proc.fill.fore_color.rgb = RGBColor(240, 253, 244)
    c_proc.line.color.rgb = BORDER_GREEN
    c_proc.line.width = Pt(1.5)
    tf_pr = c_proc.text_frame
    p = tf_pr.paragraphs[0]
    p.text = "2. THERMAL ENGINE &\nPREPROCESSING"
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = BORDER_GREEN
    p.alignment = PP_ALIGN.CENTER

    pr_items = [
        "Spatial Harmonization & Alignment",
        "Vapor Pressure & Wind Downscaling",
        "UTCI 6th-Order Polynomial (Fiala)",
        "ISO 7243 Outdoor WBGT Psychrometrics",
        "Cumulative Persistence Engine (Dmult)"
    ]
    for itm in pr_items:
        p_itm = tf_pr.add_paragraph()
        p_itm.text = f"• {itm}"
        p_itm.font.size = Pt(8.5)
        p_itm.font.color.rgb = TEXT_BLACK
        p_itm.space_before = Pt(4)

    # 3. AI RISK & VULNERABILITY MODEL
    c_risk = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(1.55), Inches(3.2), Inches(3.8))
    c_risk.fill.solid()
    c_risk.fill.fore_color.rgb = RGBColor(254, 242, 242)
    c_risk.line.color.rgb = BORDER_RED
    c_risk.line.width = Pt(1.5)
    tf_rk = c_risk.text_frame
    p = tf_rk.paragraphs[0]
    p.text = "3. AI RISK & HVI\nSYNTHESIS"
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = BORDER_RED
    p.alignment = PP_ALIGN.CENTER

    rk_items = [
        "Demographic Heat Vulnerability (HVI)",
        "Elderly & Outdoor Labor Weighting",
        "Slum Density & Canopy Deficit",
        "Normalized 0–100 Heat-Health Score",
        "Relative Risk Epidemiological Model"
    ]
    for itm in rk_items:
        p_itm = tf_rk.add_paragraph()
        p_itm.text = f"• {itm}"
        p_itm.font.size = Pt(8.5)
        p_itm.font.color.rgb = TEXT_BLACK
        p_itm.space_before = Pt(4)

    # 4. FINAL OUTPUT & ACTION
    c_out = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.0), Inches(1.55), Inches(2.533), Inches(3.8))
    c_out.fill.solid()
    c_out.fill.fore_color.rgb = BG_LIGHT_BLUE
    c_out.line.color.rgb = BORDER_BLUE
    c_out.line.width = Pt(1.5)
    tf_out = c_out.text_frame
    p = tf_out.paragraphs[0]
    p.text = "4. ACTION DELIVERY &\nVALIDATION"
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_SIH
    p.alignment = PP_ALIGN.CENTER

    out_items = [
        "Interactive Leaflet GIS Map",
        "25 REST APIs (FastAPI)",
        "Bilingual NDMA Action Playbooks",
        "Hospital Surge Bed Automation",
        "NIOSH Labor Halt Protocols"
    ]
    for itm in out_items:
        p_itm = tf_out.add_paragraph()
        p_itm.text = f"• {itm}"
        p_itm.font.size = Pt(8.5)
        p_itm.font.color.rgb = TEXT_BLACK
        p_itm.space_before = Pt(4)

    # Bottom Section: Technologies to be used (Reference Style)
    tech_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.733), Inches(1.4))
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = BG_LIGHT_GRAY
    tech_box.line.color.rgb = BORDER_GREY
    tech_box.line.width = Pt(1.2)
    tf_tb = tech_box.text_frame
    tf_tb.word_wrap = True

    p = tf_tb.paragraphs[0]
    p.text = "Technologies to be used:"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY_SIH

    p_tech = tf_tb.add_paragraph()
    p_tech.text = "• Core Backend: Python 3.11, FastAPI, Pydantic v2, Uvicorn, REST Architecture (25 Production Endpoints)\n• Scientific Computing & GIS: NumPy, Pandas, GeoPandas, Shapely, Scipy (UTCI 6th-order polynomial & psychrometrics)\n• Web GIS & Visualization: Leaflet.js (Choropleth GIS), Chart.js (5-Day Trend Curves), HTML5/CSS3 Responsive UI\n• Deployment & Data: Docker, Docker Compose, Linux VPS (Oracle Linux 6.17), Open-Meteo NWP, NASA POWER, Census India"
    p_tech.font.size = Pt(9)
    p_tech.font.color.rgb = TEXT_BLACK
    p_tech.space_before = Pt(2)

    add_reference_footer(s3, 3)

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY (EXACT REFERENCE FORMAT)
    # =========================================================================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s4)
    add_reference_header(s4, "FEASIBILITY AND VIABILITY", "ThermoShield")

    # Upper Text: Feasibility Assessment & Product Offering
    top_tb4 = s4.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.733), Inches(1.7))
    tf_t4 = top_tb4.text_frame
    tf_t4.word_wrap = True
    tf_t4.margin_left = tf_t4.margin_top = tf_t4.margin_right = tf_t4.margin_bottom = Inches(0)

    p1 = tf_t4.paragraphs[0]
    p1.text = "❖  Is this Idea feasible ?"
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = NAVY_SIH

    p1_sub = tf_t4.add_paragraph()
    p1_sub.text = "The proposed ThermoShield system is highly feasible due to the availability of open, keyless Tier-1 atmospheric datasets (Open-Meteo 0.1° NWP hourly forecasts and NASA POWER MERRA-2 40-year climatological baselines) and proven biometeorological algorithms (UTCI polynomial & ISO 7243 WBGT). The integration of Census 2011 PCA demographics enables hyper-local ward-level vulnerability mapping without waiting for restricted government clearance."
    p1_sub.font.size = Pt(9.5)
    p1_sub.font.color.rgb = TEXT_BLACK
    p1_sub.space_before = Pt(2)

    p2 = tf_t4.add_paragraph()
    p2.text = "❖  Product Offering: We are offering this system to Municipal Corporations, State Disaster Management Authorities (SDMAs), Health Departments, and Gig/Labor platforms as an automated, impact-to-action early warning and resource coordination OS."
    p2.font.size = Pt(9.5)
    p2.font.bold = True
    p2.font.color.rgb = NAVY_SIH
    p2.space_before = Pt(4)

    # Lower Left Container: Challenges & Risks vs Strategies for Overcoming (Blue Border)
    c_left = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.0), Inches(7.5), Inches(3.9))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = BG_WHITE
    c_left.line.color.rgb = BORDER_BLUE
    c_left.line.width = Pt(1.5)

    tf_cl = c_left.text_frame
    tf_cl.word_wrap = True
    tf_cl.margin_left = tf_cl.margin_top = tf_cl.margin_right = tf_cl.margin_bottom = Inches(0)

    # Column Headers inside container
    tb_heads = s4.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(7.1), Inches(0.4))
    tf_h = tb_heads.text_frame
    p_h = tf_h.paragraphs[0]
    p_h.text = "Challenges & Risks:                                             Strategies for Overcoming:"
    p_h.font.size = Pt(10.5)
    p_h.font.bold = True
    p_h.font.color.rgb = BORDER_RED

    pairs = [
        ("Data Heterogeneity (Grid vs Wards)", "Spatial Harmonization & Geometric Attribution"),
        ("Confidential Hospital Mortality Records", "Epidemiological Calibration (Azhar 2014)"),
        ("Model Overfitting / False Death Claims", "Bounded 0–100 Heat-Health Risk Index"),
        ("False Alarm Fatigue (Single Spikes)", "Compound Multi-Day Persistence Penalty (Dmult)")
    ]

    for idx, (ch, st) in enumerate(pairs):
        tb_pair = s4.shapes.add_textbox(Inches(1.0), Inches(3.6 + idx * 0.75), Inches(7.1), Inches(0.65))
        tf_p = tb_pair.text_frame
        tf_p.word_wrap = True
        tf_p.margin_left = tf_p.margin_top = tf_p.margin_right = tf_p.margin_bottom = Inches(0)
        p = tf_p.paragraphs[0]
        
        run_ch = p.add_run()
        run_ch.text = f"• {ch}\n"
        run_ch.font.size = Pt(9.5)
        run_ch.font.bold = True
        run_ch.font.color.rgb = BORDER_RED

        run_arr = p.add_run()
        run_arr.text = f"   ➔  {st}"
        run_arr.font.size = Pt(9.5)
        run_arr.font.bold = True
        run_arr.font.color.rgb = BORDER_GREEN

    # Lower Right Container: Expected Customers (Black/Navy Border)
    c_right = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(3.0), Inches(4.033), Inches(3.9))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = BG_LIGHT_GRAY
    c_right.line.color.rgb = NAVY_SIH
    c_right.line.width = Pt(1.5)

    tf_cr = c_right.text_frame
    tf_cr.word_wrap = True
    p_cr_head = tf_cr.paragraphs[0]
    p_cr_head.text = "Expected Customers & Users:"
    p_cr_head.font.size = Pt(11)
    p_cr_head.font.bold = True
    p_cr_head.font.color.rgb = NAVY_SIH

    customers = [
        ("Municipal Corporations & Smart Cities", "For ward-level water tankers, cool roofs & cooling centers."),
        ("Health Departments & Hospitals", "For heatstroke emergency ward & IV fluid prepositioning."),
        ("State Disaster Management (SDMA / NDMA)", "For automated Heat Action Plan (HAP) trigger activation."),
        ("Gig Platforms & Construction Sector", "For mandatory NIOSH labor halt (11 AM–4 PM) compliance.")
    ]
    for c_title, c_desc in customers:
        p_c = tf_cr.add_paragraph()
        p_c.text = f"➢  {c_title}"
        p_c.font.size = Pt(9.5)
        p_c.font.bold = True
        p_c.font.color.rgb = NAVY_SIH
        p_c.space_before = Pt(6)

        p_cd = tf_cr.add_paragraph()
        p_cd.text = f"    {c_desc}"
        p_cd.font.size = Pt(8.5)
        p_cd.font.color.rgb = TEXT_MUTED

    add_reference_footer(s4, 4)

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS (EXACT REFERENCE FORMAT)
    # =========================================================================
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s5)
    add_reference_header(s5, "IMPACT AND BENEFITS", "ThermoShield")

    # Top Target Audience Line
    aud_box = s5.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.733), Inches(0.8))
    tf_aud = aud_box.text_frame
    tf_aud.word_wrap = True
    tf_aud.margin_left = tf_aud.margin_top = tf_aud.margin_right = tf_aud.margin_bottom = Inches(0)
    
    p_aud = tf_aud.paragraphs[0]
    p_aud.text = "Target Audience — Municipal Commissioners, Chief Medical Officers, Disaster Response Teams (NDMA/SDRF), Gig Workers, Construction Laborers, Urban Slum Residents, General Public."
    p_aud.font.size = Pt(10)
    p_aud.font.bold = True
    p_aud.font.color.rgb = NAVY_SIH

    p_aud_sub = tf_aud.add_paragraph()
    p_aud_sub.text = "(All above stakeholders receive tailored biometeorological intelligence and automated sector-specific action protocols.)"
    p_aud_sub.font.size = Pt(9)
    p_aud_sub.font.color.rgb = TEXT_MUTED
    p_aud_sub.space_before = Pt(2)

    # Middle Banner: 72h–120h Lead-Time Timeline
    time_shape = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(11.733), Inches(1.3))
    time_shape.fill.solid()
    time_shape.fill.fore_color.rgb = BG_LIGHT_BLUE
    time_shape.line.color.rgb = BORDER_BLUE
    time_shape.line.width = Pt(1.2)
    tf_tm = time_shape.text_frame
    tf_tm.word_wrap = True
    
    p_tm_h = tf_tm.paragraphs[0]
    p_tm_h.text = "72h–120h (3–5 DAY) PRE-EMPTIVE HEAT ACTION TIMELINE:"
    p_tm_h.font.size = Pt(10)
    p_tm_h.font.bold = True
    p_tm_h.font.color.rgb = NAVY_SIH

    p_tm_b = tf_tm.add_paragraph()
    p_tm_b.text = "• 120h (D-5): NWP model flags thermal moisture surge → Early municipal advisory issued\n• 72h (D-3): Cumulative heat persistence detected → Water tankers and mobile misting units dispatched to high-risk slums\n• 48h (D-2): Hospital surge beds & ice-bath cooling centers staged; DISCOMs alerted for grid peak load\n• 24h (D-1): Mandatory NIOSH labor halts mandated for construction & delivery workers during peak hours (11 AM–4 PM)"
    p_tm_b.font.size = Pt(8.5)
    p_tm_b.font.color.rgb = TEXT_BLACK
    p_tm_b.space_before = Pt(2)

    # Three-Column Benefits Containers (Social, Economical, Environmental)
    col_w5 = Inches(3.75)
    col_gap5 = Inches(0.24)
    left_start5 = Inches(0.8)

    benefits_data = [
        ("● Social Benefits:", [
            ("Protects Vulnerable Citizens:", "Dramatically reduces heatstroke morbidity among elderly and children."),
            ("Protects Informal Laborers:", "Enforces humane work-rest cycles for delivery and construction workers."),
            ("Bilingual Public Alerts:", "Delivers plain-language vernacular guidance via SMS/WhatsApp.")
        ], BORDER_BLUE),
        ("● Economical Benefits:", [
            ("Prevents Productivity Loss:", "Reduces lost work hours through scheduled morning/evening shifts."),
            ("Optimized Resource Routing:", "Avoids wasteful blanket relief by targeting high-vulnerability wards."),
            ("Reduces Emergency Care Costs:", "Pre-empts severe multi-organ failure and ICU admissions.")
        ], BORDER_GREEN),
        ("● Environmental & Urban Benefits:", [
            ("Urban Heat Island (UHI) Mapping:", "Identifies micro-urban hot spots and green canopy deficits."),
            ("Sustainable Urban Planning:", "Guides long-term cool roof installations and tree canopy targets."),
            ("Effective Climate Resilience:", "Operationalizes national disaster mitigation at municipal ward scale.")
        ], ORANGE_SIH)
    ]

    for c_idx, (b_title, b_items, b_color) in enumerate(benefits_data):
        c_left_pos = left_start5 + c_idx * (col_w5 + col_gap5)
        card_b = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left_pos, Inches(3.35), col_w5, Inches(3.55))
        card_b.fill.solid()
        card_b.fill.fore_color.rgb = BG_WHITE
        card_b.line.color.rgb = b_color
        card_b.line.width = Pt(1.5)

        tf_cb = card_b.text_frame
        tf_cb.word_wrap = True
        tf_cb.margin_left = tf_cb.margin_top = tf_cb.margin_right = tf_cb.margin_bottom = Inches(0)

        p = tf_cb.paragraphs[0]
        p.text = b_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = b_color

        for b_h, b_d in b_items:
            p_itm = tf_cb.add_paragraph()
            p_itm.text = f"▪  {b_h} "
            p_itm.font.size = Pt(9)
            p_itm.font.bold = True
            p_itm.font.color.rgb = NAVY_SIH
            p_itm.space_before = Pt(4)

            run_bd = p_itm.add_run()
            run_bd.text = b_d
            run_bd.font.bold = False
            run_bd.font.color.rgb = TEXT_BLACK

    add_reference_footer(s5, 5)

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES (EXACT REFERENCE FORMAT)
    # =========================================================================
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(s6)
    add_reference_header(s6, "RESEARCH AND REFERENCES", "ThermoShield")

    # Content Box with ● bullets matching reference style
    ref_box = s6.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.733), Inches(5.7))
    tf_r6 = ref_box.text_frame
    tf_r6.word_wrap = True
    tf_r6.margin_left = tf_r6.margin_top = tf_r6.margin_right = tf_r6.margin_bottom = Inches(0)

    references_list = [
        ("Universal Thermal Climate Index (UTCI) Consortium (WMO / COST Action 730):", "\"Development of the Universal Thermal Climate Index (UTCI)\", Bröde, Fiala, Błażejczyk et al., Int J Biometeorol (2012). Multi-node thermophysiological polynomial."),
        ("Occupational Heat Stress Standards (ISO 7243:2017 & NIOSH):", "\"Criteria for a Recommended Standard: Occupational Exposure to Heat and Hot Environments\", CDC / NIOSH Publication No. 2016-106. Standard for WBGT thresholds & work-rest cycles."),
        ("Epidemiological Heat-Health Grounding in Indian Cities (Ahmedabad 2010):", "\"The 2010 Ahmedabad Heat Wave: Impact on Mortality and Strategy for Resilient Cities\", Azhar et al., PLoS ONE (2014). Documented 1,344 excess deaths (43.1% surge) validating humidity-heat risk."),
        ("50-Year Multi-Decadal Heat Impact in India (PNAS Reference):", "\"Increasing probability of mass-fatality heatwaves in India\", Mazdiyasni et al., Proceedings of the National Academy of Sciences (PNAS, 2017). Proving +0.5°C threshold triggers 146% rise in mass mortality events."),
        ("National Disaster & Health Policy Guidelines (NDMA & MoHFW NCDC):", "\"National Action Plan for Heat-Related Illnesses (NAP-HRI)\", National Programme on Climate Change and Human Health (NPCCHH, 2024) & NDMA National Heatwave Guidelines."),
        ("NWP & Satellite Climatology Portals:", "National Centre for Medium Range Weather Forecasting (NCMRWF) NWP Portal, IMD Geospatial Bulletins & NASA POWER MERRA-2 (LaRC) 40-Year Climatological Baseline."),
        ("Demographic Spatial Vulnerability:", "Census of India 2011 Primary Census Abstract (PCA) Ward Demographics & Open-Meteo High-Resolution Numerical Weather Prediction API."),
        ("Research Gap Addressed:", "Bridges the operational gap between macro-scale meteorological forecasts and actionable, localized public health interventions by providing an automated impact-to-action intelligence OS.")
    ]

    for idx, (r_title, r_desc) in enumerate(references_list):
        p = tf_r6.paragraphs[0] if idx == 0 else tf_r6.add_paragraph()
        p.text = f"●  {r_title} "
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = NAVY_SIH
        p.space_before = Pt(4)

        run_rd = p.add_run()
        run_rd.text = r_desc
        run_rd.font.bold = False
        run_rd.font.color.rgb = TEXT_BLACK

    add_reference_footer(s6, 6)

    # Save outputs
    out_path_1 = "/home/ubuntu/sih26083-heat-risk/docs/sih/SIH2026_IDEA_Presentation_SIH26083_Official.pptx"
    out_path_2 = "/home/ubuntu/sih26083_research/SIH26083_SIH_Official_Submission.pptx"

    os.makedirs(os.path.dirname(out_path_1), exist_ok=True)
    os.makedirs(os.path.dirname(out_path_2), exist_ok=True)

    prs.save(out_path_1)
    prs.save(out_path_2)
    print(f"Deck successfully generated matching exact reference format at:\n1. {out_path_1}\n2. {out_path_2}")

if __name__ == "__main__":
    build_sih_exact_reference_template_deck()
