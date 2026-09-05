import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_official_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Clean Dark Navy / Professional High-Contrast Palette
    NAVY_DARK = RGBColor(15, 23, 42)        # #0F172A Dark Slate
    CARD_BG = RGBColor(30, 41, 59)          # #1E293B Card Background
    CARD_BORDER = RGBColor(51, 65, 85)      # #334155 Slate Border
    TEXT_WHITE = RGBColor(248, 250, 252)    # #F8FAFC Primary Text
    TEXT_MUTED = RGBColor(148, 163, 184)    # #94A3B8 Secondary Text
    ACCENT_ORANGE = RGBColor(249, 115, 22)  # #F97316 Heat / Warning
    ACCENT_BLUE = RGBColor(56, 189, 248)    # #38BDF8 Atmosphere / Tech
    ALERT_RED = RGBColor(239, 68, 68)       # #EF4444 Danger / Risk
    ACCENT_GREEN = RGBColor(34, 197, 94)    # #22C55E Action / Health
    ACCENT_YELLOW = RGBColor(234, 179, 8)   # #EAB308 Warning

    def set_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, slide_num, category_text, title_text):
        # Header banner
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        # Category / Kicker
        p_cat = tf.paragraphs[0]
        p_cat.text = f"SMART INDIA HACKATHON 2026  |  SLIDE {slide_num} — {category_text.upper()}"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_ORANGE
        p_cat.font.name = "Arial"
        
        # Main Title
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.font.name = "Arial"
        p_title.space_before = Pt(4)

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    # =========================================================================
    # SLIDE 1: TITLE PAGE (OFFICIAL TEMPLATE ALIGNMENT)
    # =========================================================================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s1, NAVY_DARK)

    # Top SIH Banner Card
    add_card(s1, Inches(0.8), Inches(0.6), Inches(11.733), Inches(0.8), CARD_BG, ACCENT_ORANGE)
    top_box = s1.shapes.add_textbox(Inches(1.0), Inches(0.7), Inches(11.333), Inches(0.6))
    tf1_top = top_box.text_frame
    tf1_top.margin_left = tf1_top.margin_top = tf1_top.margin_right = tf1_top.margin_bottom = Inches(0)
    p1 = tf1_top.paragraphs[0]
    p1.text = "SMART INDIA HACKATHON 2026  •  OFFICIAL PROPOSAL DECK"
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_ORANGE
    p1.font.name = "Arial"
    p1.alignment = PP_ALIGN.CENTER

    # Main Hero Title
    hero_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.733), Inches(2.2))
    tf_hero = hero_box.text_frame
    tf_hero.word_wrap = True
    tf_hero.margin_left = tf_hero.margin_top = tf_hero.margin_right = tf_hero.margin_bottom = Inches(0)
    
    p_hero_main = tf_hero.paragraphs[0]
    p_hero_main.text = "ThermoShield"
    p_hero_main.font.size = Pt(40)
    p_hero_main.font.bold = True
    p_hero_main.font.color.rgb = TEXT_WHITE
    p_hero_main.font.name = "Arial"
    
    p_hero_sub = tf_hero.add_paragraph()
    p_hero_sub.text = "AI-Driven Human Heat Risk & Biometeorological Early Warning System"
    p_hero_sub.font.size = Pt(18)
    p_hero_sub.font.bold = True
    p_hero_sub.font.color.rgb = ACCENT_BLUE
    p_hero_sub.font.name = "Arial"
    p_hero_sub.space_before = Pt(6)

    p_hero_desc = tf_hero.add_paragraph()
    p_hero_desc.text = "From Weather Forecasts → Human Physiological Thermal Stress → Localized Action"
    p_hero_desc.font.size = Pt(14)
    p_hero_desc.font.color.rgb = TEXT_MUTED
    p_hero_desc.font.name = "Arial"
    p_hero_desc.space_before = Pt(4)

    # Metadata Grid (2 Cards)
    # Left Card: Problem Statement Info
    add_card(s1, Inches(0.8), Inches(4.2), Inches(5.7), Inches(2.6), CARD_BG, CARD_BORDER)
    ps_box = s1.shapes.add_textbox(Inches(1.0), Inches(4.35), Inches(5.3), Inches(2.3))
    tf_ps = ps_box.text_frame
    tf_ps.word_wrap = True
    
    p = tf_ps.paragraphs[0]
    p.text = "PROBLEM STATEMENT SPECIFICATION"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE
    
    items_ps = [
        ("Problem Statement ID:", "SIH26083"),
        ("Problem Title:", "Extreme Heatwave Early Warning & Human Thermal Stress Index"),
        ("Ministry / Org:", "Ministry of Earth Sciences (MoES)"),
        ("Department:", "National Centre for Medium Range Weather Forecasting (NCMRWF)"),
        ("Category & Theme:", "Software  |  Disaster Management")
    ]
    for label, val in items_ps:
        p = tf_ps.add_paragraph()
        p.text = f"• {label} {val}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(3)

    # Right Card: Team & Mandate Info
    add_card(s1, Inches(6.833), Inches(4.2), Inches(5.7), Inches(2.6), CARD_BG, CARD_BORDER)
    team_box = s1.shapes.add_textbox(Inches(7.033), Inches(4.35), Inches(5.3), Inches(2.3))
    tf_team = team_box.text_frame
    tf_team.word_wrap = True
    
    p = tf_team.paragraphs[0]
    p.text = "PARTICIPATION & TEAM DETAILS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    items_team = [
        ("Team Name:", "[YOUR REGISTERED TEAM NAME]"),
        ("Team ID:", "[YOUR REGISTERED TEAM ID]"),
        ("Core Mandate:", "Predict 'What weather will do to humans' vs merely 'What temperature will be'"),
        ("Key Methodology:", "Universal Thermal Climate Index (UTCI) + ISO 7243 WBGT + Census HVI"),
        ("Target Users:", "MoES/NCMRWF, NDMA, Municipal Corps, Health Depts, Gig & Labor Workers")
    ]
    for label, val in items_team:
        p = tf_team.add_paragraph()
        p.text = f"• {label} {val}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(3)

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION & IDEA
    # =========================================================================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s2, NAVY_DARK)
    add_header(s2, 2, "Idea & Innovation", "ThermoShield: Operational Human Heat-Risk Architecture")

    # Left Column: Dominant End-to-End Conceptual Flowchart
    add_card(s2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.9), CARD_BG, CARD_BORDER)
    flow_box = s2.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(5.2), Inches(4.6))
    tf_flow = flow_box.text_frame
    tf_flow.word_wrap = True
    
    p = tf_flow.paragraphs[0]
    p.text = "CORE CONCEPTUAL PIPELINE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    flow_steps = [
        ("1. Multi-Variable Ingestion", "NWP 5-day forecast + 30-yr MERRA-2 baseline (Temp, RH, Wind, Solar Rad)"),
        ("2. Human Thermal Stress Engine", "UTCI (Fiala 187-node multi-node model) + ISO 7243 WBGT + Heat Index"),
        ("3. Demographic Vulnerability (HVI)", "Census 2011 PCA spatial overlay (Elderly ≥60, Outdoor Labor, Slum Density, Green Deficit)"),
        ("4. Predictive Health-Risk Model", "Relative heat-health strain score (0–100) with multi-day accumulation penalty"),
        ("5. Automated Action Dispatcher", "Bilingual (EN/HI) NDMA playbooks, hospital surge triggers, NIOSH work-rest cycles")
    ]
    for title, desc in flow_steps:
        p = tf_flow.add_paragraph()
        p.text = f"▼ {title}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_ORANGE
        p.space_before = Pt(6)
        
        p_d = tf_flow.add_paragraph()
        p_d.text = f"   {desc}"
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED

    # Right Column: 4 Key Differentiators (Visual Cards)
    diff_cards = [
        ("1. Multidimensional Thermal Stress", "Does not rely on dry-bulb temperature alone. Accurately computes sweat evaporation limits (UTCI/WBGT) where high humidity turns 40°C lethal.", ACCENT_ORANGE),
        ("2. Human-Centric Risk vs Weather", "Translates meteorological data into actionable physiological strain and population vulnerability rather than raw atmospheric parameters.", ACCENT_BLUE),
        ("3. Hyper-Local Ward Prioritization", "Differentiates affluent shaded residential sectors from densely packed informal settlements with tin roofs and high elderly/worker shares.", ACCENT_YELLOW),
        ("4. Automated Action & Advisory Trigger", "Converts risk directly into sector-specific advisories: gig worker hydration quotas, construction halts (11 AM–4 PM), and hospital surge beds.", ACCENT_GREEN)
    ]
    
    card_w = Inches(5.9)
    card_h = Inches(1.12)
    top_start = Inches(1.5)
    for i, (title, text, color) in enumerate(diff_cards):
        y_pos = top_start + i * Inches(1.26)
        add_card(s2, Inches(6.633), y_pos, card_w, card_h, CARD_BG, color)
        tb = s2.shapes.add_textbox(Inches(6.8), y_pos + Inches(0.08), Inches(5.56), card_h - Inches(0.16))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_b = tf.add_paragraph()
        p_b.text = text
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = TEXT_WHITE
        p_b.space_before = Pt(2)

    # Bottom One-Line USP Bar
    add_card(s2, Inches(0.8), Inches(6.55), Inches(11.733), Inches(0.55), CARD_BG, ACCENT_BLUE)
    usp_box = s2.shapes.add_textbox(Inches(1.0), Inches(6.62), Inches(11.333), Inches(0.4))
    tf_usp = usp_box.text_frame
    p_usp = tf_usp.paragraphs[0]
    p_usp.text = "KEY USP: Not another heat map — an automated, end-to-end impact-to-action intelligence layer for heatwaves."
    p_usp.font.size = Pt(11)
    p_usp.font.bold = True
    p_usp.font.color.rgb = TEXT_WHITE
    p_usp.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH & ARCHITECTURE
    # =========================================================================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s3, NAVY_DARK)
    add_header(s3, 3, "Technical Approach", "5-Layer Production Architecture & Methodology")

    # 5 Horizontal Architecture Layers
    layers = [
        ("LAYER 1: DATA INGESTION", "Open-Meteo 5-Day NWP Hourly Forecast (0.1° Grid) + NASA POWER MERRA-2 30-Yr Climatological Baseline + Census 2011 PCA Data", ACCENT_BLUE),
        ("LAYER 2: HARMONIZATION & PREPROCESSING", "Coordinate Spatial Alignment + Vapor Pressure (Buck 1981) + Wind Profile Downscaling (10m → 1.2m human center-of-mass)", TEXT_MUTED),
        ("LAYER 3: THERMAL STRESS ENGINE", "Universal Thermal Climate Index (6th-order UTCI polynomial) + ISO 7243 Outdoor WBGT (Stull Tw + Liljegren Tg) + NOAA Heat Index", ACCENT_ORANGE),
        ("LAYER 4: AI & VULNERABILITY SYNTHESIS", "Demographic HVI Weighting (Elderly 25%, Labor 25%, Slum 20%, Child 15%, Canopy Deficit 15%) + Multi-Day Persistence Penalty (Dmult)", ALERT_RED),
        ("LAYER 5: DECISION & ACTION DELIVERY", "Interactive Leaflet GIS Choropleth + Chart.js Multi-Day Curves + Bilingual NDMA Playbooks + Telegram Broadcaster + 25 REST Endpoints", ACCENT_GREEN)
    ]

    layer_w = Inches(11.733)
    layer_h = Inches(0.72)
    top_l = Inches(1.45)
    for i, (l_title, l_desc, color) in enumerate(layers):
        y_l = top_l + i * Inches(0.82)
        add_card(s3, Inches(0.8), y_l, layer_w, layer_h, CARD_BG, color)
        tb = s3.shapes.add_textbox(Inches(1.0), y_l + Inches(0.06), Inches(11.333), layer_h - Inches(0.12))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = l_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_d = tf.add_paragraph()
        p_d.text = l_desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_WHITE
        p_d.space_before = Pt(2)

    # Bottom Split: Tech Stack & Operational Standards
    add_card(s3, Inches(0.8), Inches(5.65), Inches(5.7), Inches(1.45), CARD_BG, CARD_BORDER)
    ts_box = s3.shapes.add_textbox(Inches(1.0), Inches(5.75), Inches(5.3), Inches(1.25))
    tf_ts = ts_box.text_frame
    tf_ts.word_wrap = True
    p = tf_ts.paragraphs[0]
    p.text = "COMPACT & REPRODUCIBLE TECH STACK"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    t_items = [
        "Backend & API: Python 3.11, FastAPI, Pydantic v2, Uvicorn (25 REST routes)",
        "Scientific & GIS: NumPy, Pandas, GeoPandas, Shapely (Polynomial & Psychrometrics)",
        "Frontend & Viz: Leaflet.js (Choropleth GIS), Chart.js (5-day curves), Vanilla HTML5/CSS3",
        "Deployment: Docker, Docker Compose, Linux VPS (OCI Oracle 6.17), GitHub CI"
    ]
    for item in t_items:
        p = tf_ts.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(2)

    add_card(s3, Inches(6.833), Inches(5.65), Inches(5.7), Inches(1.45), CARD_BG, CARD_BORDER)
    std_box = s3.shapes.add_textbox(Inches(7.033), Inches(5.75), Inches(5.3), Inches(1.25))
    tf_std = std_box.text_frame
    tf_std.word_wrap = True
    p = tf_std.paragraphs[0]
    p.text = "STANDARDS & SCIENTIFIC CITATIONS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    s_items = [
        "UTCI Consortium: Fiala 187-node human thermoregulation polynomial (WMO standard)",
        "ISO 7243:2017 & NIOSH: Occupational heat-stress criteria and work-rest ratios",
        "NDMA 2024 & NCDC NPCCHH: National Action Plan for Heat-Related Illnesses",
        "Epidemiological Proof: Azhar et al. (2014) PLoS ONE & Mazdiyasni et al. (2017) PNAS"
    ]
    for item in s_items:
        p = tf_std.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(2)

    # =========================================================================
    # SLIDE 4: FEASIBILITY, VIABILITY & ROADMAP
    # =========================================================================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s4, NAVY_DARK)
    add_header(s4, 4, "Feasibility & Viability", "Technical Risk Mitigation & 4-Phase Deployment Roadmap")

    # 3-Column Risk & Feasibility Table (Top Section)
    col_w = Inches(3.75)
    col_gap = Inches(0.24)
    left_start = Inches(0.8)

    feasibility_data = [
        ("CHALLENGE & RISK", [
            ("Data Heterogeneity:", "Varying spatial resolutions between NWP grids (0.1°) and municipal ward polygons."),
            ("Health Data Inaccessibility:", "Hospital admission records and clinical mortality data are confidential/restricted."),
            ("Model Over-Fitting / Hallucination:", "Predicting absolute death counts without clinical validation invites severe error."),
            ("False Alarm Fatigue:", "Single temperature spikes triggering premature citywide lockdowns.")
        ], ALERT_RED),
        ("OUR ENGINEERING MITIGATION", [
            ("Spatial Harmonization:", "Dynamic geometric attribution over Census ward boundaries with population weighting."),
            ("Public Baseline Adaptation:", "Grounded in peer-reviewed epidemiological benchmarks (Azhar 2014, Mazdiyasni 2017)."),
            ("Relative Risk Indexing:", "Bounded 0–100 human heat-health risk score with NDMA-aligned action bands."),
            ("Multi-Day Persistence Penalty:", "Compound duration formula requiring sustained physiological strain to escalate alerts.")
        ], ACCENT_BLUE),
        ("WHY IT IS HIGHLY FEASIBLE", [
            ("Operational Data Ecosystem:", "Open-Meteo & NASA POWER provide open, keyless, reliable programmatic feeds today."),
            ("Tier-1 Prototype Ready:", "100% functional without government clearance, human approvals, or private API keys."),
            ("Validated Monotonicity:", "Exhaustive tests (38 test suites passing) ensure robust physical calculations."),
            ("Zero Infrastructure Cost:", "Lightweight, containerized architecture deployable on commodity VPS or cloud.")
        ], ACCENT_GREEN)
    ]

    for c_idx, (col_title, items, color) in enumerate(feasibility_data):
        c_left = left_start + c_idx * (col_w + col_gap)
        add_card(s4, c_left, Inches(1.45), col_w, Inches(3.35), CARD_BG, color)
        tb = s4.shapes.add_textbox(c_left + Inches(0.15), Inches(1.55), col_w - Inches(0.3), Inches(3.15))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = col_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        for head, body in items:
            p = tf.add_paragraph()
            p.text = f"• {head} {body}"
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT_WHITE
            p.space_before = Pt(3)

    # Bottom Section: MVP to Scale Deployment Roadmap (4 Phases)
    phases = [
        ("PHASE 1 (COMPLETED)", "Tier-1 Autonomous Prototype\n• Full UTCI / WBGT engines\n• Delhi NCR baseline\n• Keyless NWP & NASA data", ACCENT_BLUE),
        ("PHASE 2 (COMPLETED)", "Multi-City Ward Expansion\n• 5 Metros (DL, AH, ST, BB, MU)\n• Census 2011 PCA HVI overlay\n• Bilingual NDMA playbooks", ACCENT_GREEN),
        ("PHASE 3 (NEXT 3 MONTHS)", "NCMRWF Ensemble Integration\n• Direct MoES GRIB2 streams\n• MOSDAC INSAT LST integration\n• Automated WhatsApp dispatch", ACCENT_YELLOW),
        ("PHASE 4 (SCALE / PRODUCTION)", "National Health Surveillance\n• Integration with IDSP / NHRIDS\n• Machine learning calibration\n• Nationwide municipal rollout", ACCENT_ORANGE)
    ]

    p_w = Inches(2.78)
    p_gap = Inches(0.2)
    for p_idx, (p_title, p_desc, p_color) in enumerate(phases):
        p_left = left_start + p_idx * (p_w + p_gap)
        add_card(s4, p_left, Inches(4.95), p_w, Inches(2.15), CARD_BG, p_color)
        tb = s4.shapes.add_textbox(p_left + Inches(0.12), Inches(5.05), p_w - Inches(0.24), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = p_color
        
        p_b = tf.add_paragraph()
        p_b.text = p_desc
        p_b.font.size = Pt(8.5)
        p_b.font.color.rgb = TEXT_WHITE
        p_b.space_before = Pt(3)

    # =========================================================================
    # SLIDE 5: IMPACT, USER BENEFICIARIES & TIMELINE
    # =========================================================================
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s5, NAVY_DARK)
    add_header(s5, 5, "Impact & Benefits", "Target Stakeholders & 72h–120h Early Action Timeline")

    # Top Half: 72h-120h Lead-Time Timeline (5 Steps)
    add_card(s5, Inches(0.8), Inches(1.45), Inches(11.733), Inches(2.2), CARD_BG, ACCENT_BLUE)
    t_box = s5.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.333), Inches(2.0))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    
    p = tf_t.paragraphs[0]
    p.text = "PROVEN 72h–120h (3–5 DAY) PRE-EMPTIVE ACTION TIMELINE"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    timeline_points = [
        ("D-5 (120h Lead)", "NWP forecast detects moisture surge & high radiation. UTCI indicates 'Very Strong Stress' (42°C). Early municipal advisory issued.", ACCENT_YELLOW),
        ("D-3 (72h Lead)", "Cumulative heat persistence detected. High-vulnerability wards (slums, labor hubs) flagged for targeted water tanker routing & cooling centers.", ACCENT_ORANGE),
        ("D-2 (48h Lead)", "Automated bilingual bulletins sent to DISCOMs for grid load preparation and Health Departments for ORS/cooling ward staging.", ACCENT_ORANGE),
        ("D-1 (24h Lead)", "Mandatory NIOSH work-rest schedules published for gig platforms (Zomato/Swiggy) and construction sites (labor halt 11 AM–4 PM).", ALERT_RED),
        ("D-Day (Peak)", "Extreme heatwave peaks. Because interventions were pre-positioned 72h–120h prior, mass-casualty surge is successfully averted.", ACCENT_GREEN)
    ]
    for step, desc, color in timeline_points:
        p = tf_t.add_paragraph()
        p.text = f"• {step}: {desc}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(2)

    # Bottom Half: 5 Target Beneficiaries & Triple Impact
    user_groups = [
        ("Municipal Corporations", "Ward-level prioritization for water tankers, cool roofs, and emergency misting centers.", ACCENT_ORANGE),
        ("Health Departments", "Hospital surge bed activation, emergency IV fluid prepositioning, and heatstroke triage.", ALERT_RED),
        ("Disaster Management (NDMA)", "Automated activation of city Heat Action Plans (HAP) based on physiological triggers.", ACCENT_BLUE),
        ("Outdoor & Gig Workers", "Enforceable NIOSH work-rest cycles and hydration quotas for delivery, construction, and police.", ACCENT_GREEN),
        ("Citizen Public Alerts", "Bilingual localized WhatsApp/Telegram warnings enabling vulnerable families to stay safe.", ACCENT_YELLOW)
    ]

    u_w = Inches(2.22)
    u_gap = Inches(0.15)
    for u_idx, (u_title, u_desc, u_col) in enumerate(user_groups):
        u_left = left_start + u_idx * (u_w + u_gap)
        add_card(s5, u_left, Inches(3.8), u_w, Inches(3.3), CARD_BG, u_col)
        tb = s5.shapes.add_textbox(u_left + Inches(0.1), Inches(3.9), u_w - Inches(0.2), Inches(3.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = u_title
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = u_col
        
        p_b = tf.add_paragraph()
        p_b.text = u_desc
        p_b.font.size = Pt(8.5)
        p_b.font.color.rgb = TEXT_WHITE
        p_b.space_before = Pt(4)

    # =========================================================================
    # SLIDE 6: RESEARCH FOUNDATION & REFERENCES
    # =========================================================================
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s6, NAVY_DARK)
    add_header(s6, 6, "Research & References", "Authoritative Scientific Grounding & Literature Baseline")

    # 4 Research Buckets (2x2 Grid)
    grid_w = Inches(5.7)
    grid_h = Inches(2.35)
    
    buckets = [
        (Inches(0.8), Inches(1.45), "1. GOVERNMENT & METEOROLOGICAL STANDARDS", [
            ("IMD Heatwave Guidance:", "Official definitions, districtwise alerts & experimental heat index guidance."),
            ("NCMRWF NWP System:", "Target operational modeling framework and global ensemble forecast stream."),
            ("NDMA & NCDC Guidelines:", "National Action Plan for Heat-Related Illnesses (NPCCHH, 2024).")
        ], ACCENT_ORANGE),
        (Inches(6.833), Inches(1.45), "2. THERMAL STRESS & PHYSIOLOGICAL SCIENCE", [
            ("UTCI Consortium (WMO):", "Fiala multi-node thermophysiological polynomial (Bröde et al. / Jäger 2012)."),
            ("ISO 7243:2017 / NIOSH:", "Occupational exposure criteria for hot environments (CDC/NIOSH Pub 2016-106)."),
            ("Psychrometric Formulations:", "Stull (2011) natural wet bulb and Liljegren (2008) outdoor globe radiation.")
        ], ACCENT_BLUE),
        (Inches(0.8), Inches(3.95), "3. EPIDEMIOLOGICAL BENCHMARKS & VALIDATION", [
            ("Azhar et al. (2014) PLoS ONE:", "Documented 1,344 excess deaths (43.1% surge) during Ahmedabad 2010 heatwave."),
            ("Mazdiyasni et al. (2017) PNAS:", "50-year Indian analysis proving +0.5°C rise leads to 146% surge in mass mortality events."),
            ("WHO Heat & Health Evidence:", "Established empirical links between sustained thermal stress and cardiovascular collapse.")
        ], ACCENT_GREEN),
        (Inches(6.833), Inches(3.95), "4. OPEN TIER-1 DATA SOURCES & BENCHMARKS", [
            ("NASA POWER (LaRC MERRA-2):", "40-year hourly/daily climatological normals for robust baseline departure analysis."),
            ("Open-Meteo NWP Forecast:", "5-day hourly high-resolution weather variables without proprietary credentials."),
            ("Census India 2011 PCA:", "Authoritative open demographic baseline for ward-level vulnerability mapping.")
        ], ACCENT_YELLOW)
    ]

    for x, y, title, items, color in buckets:
        add_card(s6, x, y, grid_w, grid_h, CARD_BG, color)
        tb = s6.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), grid_w - Inches(0.3), grid_h - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = color
        
        for head, body in items:
            p = tf.add_paragraph()
            p.text = f"• {head} {body}"
            p.font.size = Pt(8.5)
            p.font.color.rgb = TEXT_WHITE
            p.space_before = Pt(2)

    # Bottom Callout Box: The Critical Research Gap
    add_card(s6, Inches(0.8), Inches(6.45), Inches(11.733), Inches(0.65), CARD_BG, ACCENT_BLUE)
    gap_box = s6.shapes.add_textbox(Inches(1.0), Inches(6.52), Inches(11.333), Inches(0.5))
    tf_gap = gap_box.text_frame
    p_gap = tf_gap.paragraphs[0]
    p_gap.text = "RESEARCH GAP SOLVED: Existing systems provide either macro-weather forecasts OR static vulnerability maps. ThermoShield operationalizes the complete unified pipeline: Forecast → Physiological Thermal Stress → Demographic Vulnerability → Ward Risk → Automated Action."
    p_gap.font.size = Pt(9.5)
    p_gap.font.bold = True
    p_gap.font.color.rgb = TEXT_WHITE
    p_gap.alignment = PP_ALIGN.CENTER

    # Save to both target locations
    target_path_1 = "/home/ubuntu/sih26083-heat-risk/docs/sih/ThermoShield_SIH26083_Official_Presentation.pptx"
    target_path_2 = "/home/ubuntu/sih26083_research/SIH26083_SIH_Official_Submission.pptx"
    
    os.makedirs(os.path.dirname(target_path_1), exist_ok=True)
    os.makedirs(os.path.dirname(target_path_2), exist_ok=True)
    
    prs.save(target_path_1)
    prs.save(target_path_2)
    print(f"Presentation saved successfully to:\n1. {target_path_1}\n2. {target_path_2}")

if __name__ == "__main__":
    build_official_deck()
