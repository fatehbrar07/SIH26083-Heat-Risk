import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_sih2026_official_template_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Official SIH Brand & Government Color Palette (Crisp Light Theme)
    BG_LIGHT = RGBColor(248, 250, 252)       # #F8FAFC Soft Off-White Background
    CARD_BG = RGBColor(255, 255, 255)        # #FFFFFF Pure White Cards
    CARD_BORDER = RGBColor(226, 232, 240)    # #E2E8F0 Subtle Border
    
    NAVY_PRIMARY = RGBColor(15, 44, 89)      # #0F2C59 Official SIH Deep Navy
    ORANGE_ACCENT = RGBColor(234, 88, 12)    # #EA580C SIH Saffron / Orange
    TEXT_DARK = RGBColor(30, 41, 59)         # #1E293B Primary Text
    TEXT_MUTED = RGBColor(100, 116, 139)     # #64748B Secondary Text
    
    ALERT_RED = RGBColor(220, 38, 38)        # #DC2626 Critical Alert
    GREEN_SUCCESS = RGBColor(22, 163, 74)    # #16A34A Verified / Action
    BLUE_ACCENT = RGBColor(2, 132, 199)      # #0284C7 Atmosphere / Tech

    def set_bg(slide, color=BG_LIGHT):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.2)
        else:
            shape.line.fill.background()
        return shape

    def add_sih_header(slide, slide_num, slide_heading, sub_heading):
        # Top Official Header Banner
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = NAVY_PRIMARY
        top_bar.line.fill.background()

        # Orange Accent Line under Top Bar
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.06))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = ORANGE_ACCENT
        accent_line.line.fill.background()

        # Header Text Box
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.733), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        # SIH Brand Kicker
        p0 = tf.paragraphs[0]
        p0.text = f"SMART INDIA HACKATHON 2026  •  OFFICIAL IDEA PRESENTATION  •  SLIDE {slide_num}/6"
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = ORANGE_ACCENT
        p0.font.name = "Arial"

        # Main Slide Title
        p1 = tf.add_paragraph()
        p1.text = slide_heading.upper()
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        p1.font.name = "Arial"
        p1.space_before = Pt(2)

        # Bottom Sub-header in slide body
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.28), Inches(11.733), Inches(0.4))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = Inches(0)
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = sub_heading
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = NAVY_PRIMARY
        p_sub.font.name = "Arial"

    # =========================================================================
    # SLIDE 1: TITLE PAGE (OFFICIAL SIH FORMAT)
    # =========================================================================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s1, BG_LIGHT)

    # Top SIH Official Header Band
    top_band = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = NAVY_PRIMARY
    top_band.line.fill.background()

    orange_stripe = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.3), Inches(13.333), Inches(0.08))
    orange_stripe.fill.solid()
    orange_stripe.fill.fore_color.rgb = ORANGE_ACCENT
    orange_stripe.line.fill.background()

    # Title in Header
    t_box = s1.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.733), Inches(0.9))
    tf_t = t_box.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.text = "SMART INDIA HACKATHON 2026"
    p_t.font.size = Pt(24)
    p_t.font.bold = True
    p_t.font.color.rgb = RGBColor(255, 255, 255)
    p_t.alignment = PP_ALIGN.CENTER
    
    p_sub_t = tf_t.add_paragraph()
    p_sub_t.text = "Official Idea Presentation Template  |  Ministry of Earth Sciences (MoES)"
    p_sub_t.font.size = Pt(12)
    p_sub_t.font.color.rgb = ORANGE_ACCENT
    p_sub_t.alignment = PP_ALIGN.CENTER
    p_sub_t.space_before = Pt(2)

    # Solution Hero Card
    add_card(s1, Inches(0.8), Inches(1.6), Inches(11.733), Inches(1.8), CARD_BG, ORANGE_ACCENT)
    hero_tb = s1.shapes.add_textbox(Inches(1.1), Inches(1.75), Inches(11.133), Inches(1.5))
    tf_h = hero_tb.text_frame
    tf_h.word_wrap = True
    tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = Inches(0)
    
    p_sol_kicker = tf_h.paragraphs[0]
    p_sol_kicker.text = "PROPOSED SOLUTION NAME"
    p_sol_kicker.font.size = Pt(11)
    p_sol_kicker.font.bold = True
    p_sol_kicker.font.color.rgb = ORANGE_ACCENT
    
    p_sol_title = tf_h.add_paragraph()
    p_sol_title.text = "ThermoShield — AI-Driven Human Heat Risk & Early Warning System"
    p_sol_title.font.size = Pt(24)
    p_sol_title.font.bold = True
    p_sol_title.font.color.rgb = NAVY_PRIMARY
    p_sol_title.space_before = Pt(2)

    p_sol_sub = tf_h.add_paragraph()
    p_sol_sub.text = "From Weather Forecasts → Physiological Human Thermal Stress → Hyper-Local Ward Action"
    p_sol_sub.font.size = Pt(13)
    p_sol_sub.font.bold = True
    p_sol_sub.font.color.rgb = BLUE_ACCENT
    p_sol_sub.space_before = Pt(4)

    # Metadata Two-Column Grid (Official Fields)
    # Left Card: Problem Statement Metadata
    add_card(s1, Inches(0.8), Inches(3.6), Inches(5.7), Inches(3.4), CARD_BG, CARD_BORDER)
    ps_tb = s1.shapes.add_textbox(Inches(1.0), Inches(3.75), Inches(5.3), Inches(3.1))
    tf_ps = ps_tb.text_frame
    tf_ps.word_wrap = True
    tf_ps.margin_left = tf_ps.margin_top = tf_ps.margin_right = tf_ps.margin_bottom = Inches(0)
    
    p = tf_ps.paragraphs[0]
    p.text = "PROBLEM STATEMENT DETAILS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY
    
    ps_fields = [
        ("Problem Statement ID:", "SIH26083"),
        ("Problem Statement Title:", "Extreme Heatwave Early Warning and Human Thermal Stress Index"),
        ("Organization:", "Ministry of Earth Sciences (MoES)"),
        ("Department:", "National Centre for Medium Range Weather Forecasting (NCMRWF)"),
        ("Theme:", "Disaster Management"),
        ("PS Category:", "Software")
    ]
    for label, val in ps_fields:
        p = tf_ps.add_paragraph()
        p.text = f"• {label} "
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = ORANGE_ACCENT
        p.space_before = Pt(4)
        
        run_v = p.add_run()
        run_v.text = val
        run_v.font.bold = False
        run_v.font.color.rgb = TEXT_DARK

    # Right Card: Team & Registration Metadata
    add_card(s1, Inches(6.833), Inches(3.6), Inches(5.7), Inches(3.4), CARD_BG, CARD_BORDER)
    team_tb = s1.shapes.add_textbox(Inches(7.033), Inches(3.75), Inches(5.3), Inches(3.1))
    tf_team = team_tb.text_frame
    tf_team.word_wrap = True
    tf_team.margin_left = tf_team.margin_top = tf_team.margin_right = tf_team.margin_bottom = Inches(0)
    
    p = tf_team.paragraphs[0]
    p.text = "TEAM & SUBMISSION DETAILS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY
    
    team_fields = [
        ("Team ID:", "[YOUR REGISTERED TEAM ID]"),
        ("Team Name:", "[YOUR REGISTERED TEAM NAME]"),
        ("Team Leader:", "[TEAM LEADER NAME]"),
        ("Institute Name:", "[YOUR COLLEGE / UNIVERSITY NAME]"),
        ("Core Science:", "Universal Thermal Climate Index (UTCI) + ISO 7243 WBGT"),
        ("Prototype Status:", "Fully runnable Tier-1 prototype with multi-city GIS & REST API")
    ]
    for label, val in team_fields:
        p = tf_team.add_paragraph()
        p.text = f"• {label} "
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = GREEN_SUCCESS
        p.space_before = Pt(4)
        
        run_v = p.add_run()
        run_v.text = val
        run_v.font.bold = False
        run_v.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION & IDEA TITLE
    # =========================================================================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s2, BG_LIGHT)
    add_sih_header(s2, 2, "Idea & Proposed Solution", "ThermoShield: Unified Impact-Based Heat Risk Operating System")

    # Left Column: Dominant End-to-End Conceptual Flowchart
    add_card(s2, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.95), CARD_BG, CARD_BORDER)
    flow_tb = s2.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.2), Inches(4.65))
    tf_f = flow_tb.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = Inches(0)
    
    p = tf_f.paragraphs[0]
    p.text = "END-TO-END OPERATIONAL FLOWCHART"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY
    
    f_nodes = [
        ("WEATHER + SATELLITE + DEMOGRAPHICS", "5-Day NWP Forecasts + NASA 30-Yr Climatology + Census 2011 PCA Data", BLUE_ACCENT),
        ("HUMAN THERMAL STRESS ENGINE", "UTCI (Fiala 187-Node Model) + ISO 7243 WBGT + NOAA Heat Index", ORANGE_ACCENT),
        ("AI & POPULATION RISK MODEL", "Relative Heat-Health Score (0–100) with Multi-Day Accumulation Penalty", ALERT_RED),
        ("SPATIAL VULNERABILITY (GIS)", "Ward/Zone Risk Mapping (Elderly, Outdoor Labor, Slums, Green Deficit)", NAVY_PRIMARY),
        ("AUTOMATED ACTION ENGINE", "Bilingual NDMA Advisories, Hospital Surge Triggers, NIOSH Labor Cycles", GREEN_SUCCESS)
    ]
    for idx, (n_title, n_desc, n_col) in enumerate(f_nodes):
        p = tf_f.add_paragraph()
        p.text = f"{'▼ ' if idx > 0 else '● '}{n_title}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = n_col
        p.space_before = Pt(5)
        
        p_sub = tf_f.add_paragraph()
        p_sub.text = f"   {n_desc}"
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = TEXT_MUTED

    # Right Column: 4 Key Differentiator Cards
    diff_cards = [
        ("1. Multidimensional Thermal Stress", "Does not rely on dry-bulb temperature alone. Accurately evaluates sweat evaporation breakdown (UTCI/WBGT) where high humidity turns 40°C lethal.", ORANGE_ACCENT),
        ("2. Human-Centric Risk vs Raw Weather", "Translates atmospheric parameters into physiological thermoregulatory strain and population health vulnerability.", BLUE_ACCENT),
        ("3. Hyper-Local Ward Prioritization", "Differentiates affluent shaded residential sectors from densely packed informal settlements with tin roofs and high worker concentrations.", ALERT_RED),
        ("4. Automated Action & Advisory Trigger", "Directly translates risk into sector-specific advisories: gig worker hydration quotas, construction halts (11 AM–4 PM), and hospital surge beds.", GREEN_SUCCESS)
    ]

    top_d = Inches(1.7)
    c_w = Inches(5.9)
    c_h = Inches(1.15)
    for i, (title, desc, color) in enumerate(diff_cards):
        y = top_d + i * Inches(1.27)
        add_card(s2, Inches(6.633), y, c_w, c_h, CARD_BG, color)
        tb = s2.shapes.add_textbox(Inches(6.8), y + Inches(0.08), Inches(5.56), c_h - Inches(0.16))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_DARK
        p_d.space_before = Pt(2)

    # Bottom USP Callout
    add_card(s2, Inches(0.8), Inches(6.75), Inches(11.733), Inches(0.5), CARD_BG, BLUE_ACCENT)
    usp_tb = s2.shapes.add_textbox(Inches(1.0), Inches(6.8), Inches(11.333), Inches(0.4))
    tf_u = usp_tb.text_frame
    p_u = tf_u.paragraphs[0]
    p_u.text = "ONE-LINE USP: Not another heat map — an automated, end-to-end impact-to-action intelligence layer for heatwaves."
    p_u.font.size = Pt(10.5)
    p_u.font.bold = True
    p_u.font.color.rgb = NAVY_PRIMARY
    p_u.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH & ARCHITECTURE
    # =========================================================================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s3, BG_LIGHT)
    add_sih_header(s3, 3, "Technical Approach", "5-Layer Production Architecture, Standards & Stack")

    # 5 Horizontal Architecture Layers
    layers = [
        ("LAYER 1: DATA INGESTION", "Open-Meteo 5-Day NWP Hourly Forecast (0.1° Grid) + NASA POWER MERRA-2 30-Yr Climatological Baseline + Census 2011 PCA Data", BLUE_ACCENT),
        ("LAYER 2: DATA PREPROCESSING & ALIGNMENT", "Coordinate Spatial Harmonization + Vapor Pressure (Buck 1981) + Logarithmic Wind Profile Downscaling (10m → 1.2m human center)", TEXT_MUTED),
        ("LAYER 3: THERMAL STRESS ENGINE", "Universal Thermal Climate Index (6th-order UTCI polynomial) + ISO 7243 Outdoor WBGT (Stull Tw + Liljegren Tg) + NOAA Heat Index", ORANGE_ACCENT),
        ("LAYER 4: AI & VULNERABILITY SYNTHESIS", "Demographic HVI Weighting (Elderly 25%, Labor 25%, Slum 20%, Child 15%, Canopy Deficit 15%) + Multi-Day Persistence Penalty (Dmult)", ALERT_RED),
        ("LAYER 5: DECISION & ACTION DELIVERY", "Interactive Leaflet GIS Choropleth + Chart.js Multi-Day Curves + Bilingual NDMA Playbooks + Telegram Broadcaster + 25 REST Endpoints", GREEN_SUCCESS)
    ]

    top_l = Inches(1.7)
    l_w = Inches(11.733)
    l_h = Inches(0.72)
    for i, (l_title, l_desc, color) in enumerate(layers):
        y = top_l + i * Inches(0.8)
        add_card(s3, Inches(0.8), y, l_w, l_h, CARD_BG, color)
        tb = s3.shapes.add_textbox(Inches(1.0), y + Inches(0.06), Inches(11.333), l_h - Inches(0.12))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = l_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_d = tf.add_paragraph()
        p_d.text = l_desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_DARK
        p_d.space_before = Pt(2)

    # Bottom Split: Technology Stack vs Scientific Standards
    add_card(s3, Inches(0.8), Inches(5.8), Inches(5.7), Inches(1.45), CARD_BG, CARD_BORDER)
    ts_tb = s3.shapes.add_textbox(Inches(1.0), Inches(5.9), Inches(5.3), Inches(1.25))
    tf_ts = ts_tb.text_frame
    tf_ts.word_wrap = True
    tf_ts.margin_left = tf_ts.margin_top = tf_ts.margin_right = tf_ts.margin_bottom = Inches(0)
    p = tf_ts.paragraphs[0]
    p.text = "COMPACT & REPRODUCIBLE TECH STACK"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY
    
    t_bullets = [
        "Backend & API: Python 3.11, FastAPI, Pydantic v2, Uvicorn (25 REST routes)",
        "Scientific & GIS: NumPy, Pandas, GeoPandas, Shapely (Polynomial & Psychrometrics)",
        "Frontend & Viz: Leaflet.js (Choropleth GIS), Chart.js (5-day curves), Vanilla HTML5/CSS3",
        "Deployment: Docker, Docker Compose, Linux VPS (OCI Oracle 6.17), GitHub CI"
    ]
    for b in t_bullets:
        p = tf_ts.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(2)

    add_card(s3, Inches(6.833), Inches(5.8), Inches(5.7), Inches(1.45), CARD_BG, CARD_BORDER)
    std_tb = s3.shapes.add_textbox(Inches(7.033), Inches(5.9), Inches(5.3), Inches(1.25))
    tf_std = std_tb.text_frame
    tf_std.word_wrap = True
    tf_std.margin_left = tf_std.margin_top = tf_std.margin_right = tf_std.margin_bottom = Inches(0)
    p = tf_std.paragraphs[0]
    p.text = "SCIENTIFIC CITATIONS & METHODOLOGY"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY
    
    s_bullets = [
        "UTCI Consortium: Fiala 187-node human thermoregulation polynomial (WMO standard)",
        "ISO 7243:2017 & NIOSH: Occupational heat-stress criteria and work-rest ratios",
        "NDMA 2024 & NCDC NPCCHH: National Action Plan for Heat-Related Illnesses",
        "Epidemiological Proof: Azhar et al. (2014) PLoS ONE & Mazdiyasni et al. (2017) PNAS"
    ]
    for b in s_bullets:
        p = tf_std.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(2)

    # =========================================================================
    # SLIDE 4: FEASIBILITY & VIABILITY
    # =========================================================================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s4, BG_LIGHT)
    add_sih_header(s4, 4, "Feasibility & Viability", "Technical Risk Mitigation Matrix & 4-Phase Deployment Roadmap")

    # 3-Column Risk & Mitigation Table
    col_w = Inches(3.75)
    col_gap = Inches(0.24)
    left_start = Inches(0.8)

    feas_data = [
        ("POTENTIAL CHALLENGE & RISK", [
            ("Data Heterogeneity:", "Varying spatial resolutions between NWP grids (0.1°) and municipal ward polygons."),
            ("Health Data Inaccessibility:", "Hospital admission records and clinical mortality data are confidential/restricted."),
            ("Model Over-Fitting / Hallucination:", "Predicting absolute death counts without clinical validation invites severe error."),
            ("False Alarm Fatigue:", "Single temperature spikes triggering premature citywide lockdowns.")
        ], ALERT_RED),
        ("OUR ENGINEERING MITIGATION", [
            ("Spatial Harmonization:", "Dynamic geometric attribution over Census ward boundaries with population weighting."),
            ("Public Baseline Adaptation:", "Grounded in peer-reviewed epidemiological benchmarks (Azhar 2014, Mazdiyasni 2017)."),
            ("Relative Risk Indexing:", "Bounded 0–100 human heat-health risk score with NDMA-aligned action bands."),
            ("Multi-Day Persistence Penalty:", "Compound duration formula requiring sustained physiological strain to escalate alerts.")
        ], BLUE_ACCENT),
        ("WHY IT IS HIGHLY FEASIBLE", [
            ("Operational Data Ecosystem:", "Open-Meteo & NASA POWER provide open, keyless, reliable programmatic feeds today."),
            ("Tier-1 Prototype Ready:", "100% functional without government clearance, human approvals, or private API keys."),
            ("Validated Monotonicity:", "Exhaustive tests (38 test suites passing) ensure robust physical calculations."),
            ("Zero Infrastructure Cost:", "Lightweight, containerized architecture deployable on commodity VPS or cloud.")
        ], GREEN_SUCCESS)
    ]

    for c_idx, (col_title, items, color) in enumerate(feas_data):
        c_left = left_start + c_idx * (col_w + col_gap)
        add_card(s4, c_left, Inches(1.7), col_w, Inches(3.2), CARD_BG, color)
        tb = s4.shapes.add_textbox(c_left + Inches(0.15), Inches(1.8), col_w - Inches(0.3), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = col_title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        for head, body in items:
            p = tf.add_paragraph()
            p.text = f"• {head} "
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = NAVY_PRIMARY
            p.space_before = Pt(3)
            
            run_b = p.add_run()
            run_b.text = body
            run_b.font.bold = False
            run_b.font.color.rgb = TEXT_DARK

    # Bottom Section: MVP to Scale Deployment Roadmap (4 Phases)
    phases = [
        ("PHASE 1 (COMPLETED)", "Tier-1 Autonomous Prototype\n• Full UTCI / WBGT engines\n• Delhi NCR baseline\n• Keyless NWP & NASA data", BLUE_ACCENT),
        ("PHASE 2 (COMPLETED)", "Multi-City Ward Expansion\n• 5 Metros (DL, AH, ST, BB, MU)\n• Census 2011 PCA HVI overlay\n• Bilingual NDMA playbooks", GREEN_SUCCESS),
        ("PHASE 3 (NEXT 3 MONTHS)", "NCMRWF Ensemble Integration\n• Direct MoES GRIB2 streams\n• MOSDAC INSAT LST integration\n• Automated WhatsApp dispatch", ORANGE_ACCENT),
        ("PHASE 4 (SCALE / PRODUCTION)", "National Health Surveillance\n• Integration with IDSP / NHRIDS\n• Machine learning calibration\n• Nationwide municipal rollout", ALERT_RED)
    ]

    p_w = Inches(2.78)
    p_gap = Inches(0.2)
    for p_idx, (p_title, p_desc, p_color) in enumerate(phases):
        p_left = left_start + p_idx * (p_w + p_gap)
        add_card(s4, p_left, Inches(5.05), p_w, Inches(2.15), CARD_BG, p_color)
        tb = s4.shapes.add_textbox(p_left + Inches(0.12), Inches(5.15), p_w - Inches(0.24), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = p_color
        
        p_b = tf.add_paragraph()
        p_b.text = p_desc
        p_b.font.size = Pt(8.5)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_before = Pt(3)

    # =========================================================================
    # SLIDE 5: IMPACT & BENEFITS
    # =========================================================================
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s5, BG_LIGHT)
    add_sih_header(s5, 5, "Impact & Benefits", "Target Stakeholders & 72h–120h Early Intervention Timeline")

    # Top Half: 72h-120h Lead-Time Timeline
    add_card(s5, Inches(0.8), Inches(1.7), Inches(11.733), Inches(2.2), CARD_BG, BLUE_ACCENT)
    t_tb = s5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.0))
    tf_t = t_tb.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = Inches(0)
    
    p = tf_t.paragraphs[0]
    p.text = "PROVEN 72h–120h (3–5 DAY) PRE-EMPTIVE ACTION TIMELINE"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    timeline_points = [
        ("D-5 (120h Lead)", "NWP forecast detects moisture surge & high radiation. UTCI indicates 'Very Strong Stress' (42°C). Early municipal advisory issued.", ORANGE_ACCENT),
        ("D-3 (72h Lead)", "Cumulative heat persistence detected. High-vulnerability wards (slums, labor hubs) flagged for targeted water tanker routing & cooling centers.", ORANGE_ACCENT),
        ("D-2 (48h Lead)", "Automated bilingual bulletins sent to DISCOMs for grid load preparation and Health Departments for ORS/cooling ward staging.", ALERT_RED),
        ("D-1 (24h Lead)", "Mandatory NIOSH work-rest schedules published for gig platforms (Zomato/Swiggy) and construction sites (labor halt 11 AM–4 PM).", ALERT_RED),
        ("D-Day (Peak)", "Extreme heatwave peaks. Because interventions were pre-positioned 72h–120h prior, mass-casualty surge is successfully averted.", GREEN_SUCCESS)
    ]
    for step, desc, color in timeline_points:
        p = tf_t.add_paragraph()
        p.text = f"• {step}: "
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_before = Pt(2)
        
        run_d = p.add_run()
        run_d.text = desc
        run_d.font.bold = False
        run_d.font.color.rgb = TEXT_DARK

    # Bottom Half: 5 Target Beneficiaries & Triple Impact
    user_groups = [
        ("Municipal Corporations", "Ward-level prioritization for water tankers, cool roofs, and emergency misting centers.", ORANGE_ACCENT),
        ("Health Departments", "Hospital surge bed activation, emergency IV fluid prepositioning, and heatstroke triage.", ALERT_RED),
        ("Disaster Management (NDMA)", "Automated activation of city Heat Action Plans (HAP) based on physiological triggers.", BLUE_ACCENT),
        ("Outdoor & Gig Workers", "Enforceable NIOSH work-rest cycles and hydration quotas for delivery, construction, and police.", GREEN_SUCCESS),
        ("Citizen Public Alerts", "Bilingual localized WhatsApp/Telegram warnings enabling vulnerable families to stay safe.", NAVY_PRIMARY)
    ]

    u_w = Inches(2.22)
    u_gap = Inches(0.15)
    for u_idx, (u_title, u_desc, u_col) in enumerate(user_groups):
        u_left = left_start + u_idx * (u_w + u_gap)
        add_card(s5, u_left, Inches(4.05), u_w, Inches(3.15), CARD_BG, u_col)
        tb = s5.shapes.add_textbox(u_left + Inches(0.1), Inches(4.15), u_w - Inches(0.2), Inches(2.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = u_title
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = u_col
        
        p_b = tf.add_paragraph()
        p_b.text = u_desc
        p_b.font.size = Pt(8.5)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_before = Pt(4)

    # =========================================================================
    # SLIDE 6: RESEARCH & REFERENCES
    # =========================================================================
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s6, BG_LIGHT)
    add_sih_header(s6, 6, "Research & References", "Authoritative Scientific Grounding & Literature Baseline")

    # 4 Research Buckets (2x2 Grid)
    grid_w = Inches(5.7)
    grid_h = Inches(2.25)
    
    buckets = [
        (Inches(0.8), Inches(1.7), "1. GOVERNMENT & METEOROLOGICAL STANDARDS", [
            ("IMD Heatwave Guidance:", "Official definitions, districtwise alerts & experimental heat index guidance."),
            ("NCMRWF NWP System:", "Target operational modeling framework and global ensemble forecast stream."),
            ("NDMA & NCDC Guidelines:", "National Action Plan for Heat-Related Illnesses (NPCCHH, 2024).")
        ], ORANGE_ACCENT),
        (Inches(6.833), Inches(1.7), "2. THERMAL STRESS & PHYSIOLOGICAL SCIENCE", [
            ("UTCI Consortium (WMO):", "Fiala multi-node thermophysiological polynomial (Bröde et al. / Jäger 2012)."),
            ("ISO 7243:2017 / NIOSH:", "Occupational exposure criteria for hot environments (CDC/NIOSH Pub 2016-106)."),
            ("Psychrometric Formulations:", "Stull (2011) natural wet bulb and Liljegren (2008) outdoor globe radiation.")
        ], BLUE_ACCENT),
        (Inches(0.8), Inches(4.1), "3. EPIDEMIOLOGICAL BENCHMARKS & VALIDATION", [
            ("Azhar et al. (2014) PLoS ONE:", "Documented 1,344 excess deaths (43.1% surge) during Ahmedabad 2010 heatwave."),
            ("Mazdiyasni et al. (2017) PNAS:", "50-year Indian analysis proving +0.5°C rise leads to 146% surge in mass mortality events."),
            ("WHO Heat & Health Evidence:", "Established empirical links between sustained thermal stress and cardiovascular collapse.")
        ], GREEN_SUCCESS),
        (Inches(6.833), Inches(4.1), "4. OPEN TIER-1 DATA SOURCES & BENCHMARKS", [
            ("NASA POWER (LaRC MERRA-2):", "40-year hourly/daily climatological normals for robust baseline departure analysis."),
            ("Open-Meteo NWP Forecast:", "5-day hourly high-resolution weather variables without proprietary credentials."),
            ("Census India 2011 PCA:", "Authoritative open demographic baseline for ward-level vulnerability mapping.")
        ], ALERT_RED)
    ]

    for x, y, title, items, color in buckets:
        add_card(s6, x, y, grid_w, grid_h, CARD_BG, color)
        tb = s6.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), grid_w - Inches(0.3), grid_h - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = color
        
        for head, body in items:
            p = tf.add_paragraph()
            p.text = f"• {head} "
            p.font.size = Pt(8.5)
            p.font.bold = True
            p.font.color.rgb = NAVY_PRIMARY
            p.space_before = Pt(2)
            
            run_b = p.add_run()
            run_b.text = body
            run_b.font.bold = False
            run_b.font.color.rgb = TEXT_DARK

    # Bottom Callout Box: The Critical Research Gap
    add_card(s6, Inches(0.8), Inches(6.5), Inches(11.733), Inches(0.65), CARD_BG, BLUE_ACCENT)
    gap_tb = s6.shapes.add_textbox(Inches(1.0), Inches(6.55), Inches(11.333), Inches(0.55))
    tf_gap = gap_tb.text_frame
    tf_gap.word_wrap = True
    tf_gap.margin_left = tf_gap.margin_top = tf_gap.margin_right = tf_gap.margin_bottom = Inches(0)
    p_gap = tf_gap.paragraphs[0]
    p_gap.text = "RESEARCH GAP SOLVED: Existing systems provide either macro-weather forecasts OR static vulnerability maps. ThermoShield operationalizes the complete unified pipeline: Forecast → Physiological Thermal Stress → Demographic Vulnerability → Ward Risk → Automated Action."
    p_gap.font.size = Pt(9.5)
    p_gap.font.bold = True
    p_gap.font.color.rgb = NAVY_PRIMARY
    p_gap.alignment = PP_ALIGN.CENTER

    # Save to canonical target locations
    target_path_1 = "/home/ubuntu/sih26083-heat-risk/docs/sih/SIH2026_IDEA_Presentation_SIH26083_Official.pptx"
    target_path_2 = "/home/ubuntu/sih26083_research/SIH26083_SIH_Official_Submission.pptx"
    
    os.makedirs(os.path.dirname(target_path_1), exist_ok=True)
    os.makedirs(os.path.dirname(target_path_2), exist_ok=True)
    
    prs.save(target_path_1)
    prs.save(target_path_2)
    print(f"Official SIH 2026 Template Deck successfully built and saved to:\n1. {target_path_1}\n2. {target_path_2}")

if __name__ == "__main__":
    build_sih2026_official_template_deck()
